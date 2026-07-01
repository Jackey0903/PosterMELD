from copy import deepcopy
import json
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE

from src.agents.affiliation_logo_agent import AffiliationLogoAgent
from src.agents.adaptive_column_relayout import AdaptiveColumnRelayoutAgent
from src.agents.background_image_agent import BackgroundImageAgent
from src.agents.block_content_refiner import BlockContentRefiner
from src.agents.block_occupancy_analyzer import BlockOccupancyAnalyzer
from src.agents.block_vlm_reviewer import BlockVLMReviewer
from src.agents.curator import StoryBoardCurator
from src.agents.font_agent import FontAgent
from src.agents.layout_agent import LayoutAgent
from src.agents.micro_layout_refiner import MicroLayoutRefiner
from src.agents.parser import Parser
from src.agents.poster_keypoint_selector import PosterKeypointSelector
from src.agents.renderer import Renderer
from src.agents.section_title_designer import SectionTitleDesigner
from src.agents.standard_template_preselector import StandardTemplatePreselector
from src.agents.template_capacity_planner import TemplateCapacityPlanner
from src.agents.template_block_planner import TemplateBlockPlanner
from src.agents.vlm_layout_reviewer import VLMLayoutReviewer
from src.agents.visual_asset_agent import VisualAssetAgent
from src.agents.visual_legibility_reviewer import VisualLegibilityReviewer
from src.config.poster_config import load_config
from src.layout.template_selector import TemplateSelector
from src.state.poster_state import create_state, _get_model_config
from src.template_extraction.block_template_registry import (
    build_runtime_template,
    get_block_template_info,
    list_block_template_ids,
    load_block_template_layout,
)
from src.template_extraction.extract_templates import build_template
from src.template_extraction.registry import list_extracted_template_ids, load_extracted_template
from src.tools.image_api import ImageTools
from src.tools.layout_api import LayoutTemplates
from src.tools.pptx_api import PPTXDirector
from src.utils.text_cleanup import normalize_text_for_poster, normalize_title_for_poster
from src.workflow.pipeline import _run_final_quality_gate, resolve_poster_dimensions


def test_parser_visual_assets_registry_matches_images_tables():
    parser = Parser.__new__(Parser)
    figures = {
        "1": {
            "caption": "Figure 1",
            "path": "/tmp/figure-1.png",
            "aspect": 1.5,
        }
    }
    tables = {
        "2": {
            "caption": "Table 2",
            "path": "/tmp/table-2.png",
            "aspect": 2.0,
        }
    }

    visual_assets = parser._build_visual_registry(figures, tables)

    assert visual_assets["figure_1"]["source_path"] == figures["1"]["path"]
    assert visual_assets["figure_1"]["asset_type"] == "figure"
    assert visual_assets["table_2"]["source_path"] == tables["2"]["path"]
    assert visual_assets["table_2"]["asset_type"] == "table"


def test_parser_extracts_affiliations_from_paper_header():
    parser = Parser.__new__(Parser)
    raw_text = """
    # Example Paper
    A. One, B. Two
    Department of Computer Science and Engineering, Washington University in St. Louis, USA
    Department of Computer Science and Engineering, George Mason University
    Brown School at Washington University in St. Louis, USA

    #### Abstract
    Body starts here.
    """

    affiliations = parser._extract_affiliations(raw_text)

    assert "Washington University in St. Louis" in affiliations
    assert "George Mason University" in affiliations
    assert "Brown School at Washington University in St. Louis" in affiliations


def test_parser_does_not_extract_reference_doi_as_paper_doi():
    parser = Parser.__new__(Parser)
    raw_text = """
    # Can Watermarked LLMs Be Identified By Users Via Crafted Prompts?
    Tsinghua University

    A BSTRACT
    Body starts here.

    References
    Some unrelated paper. DOI: 10.1145/3626772.3661377
    """

    assert parser._extract_doi(raw_text) is None


def test_parser_extracts_watermark_paper_affiliations_from_header():
    parser = Parser.__new__(Parser)
    raw_text = """
    # Can Watermarked LLMs Be Identified By Users Via Crafted Prompts?
    Tsinghua University
    Beijing University of Posts and Telecommunications
    The Chinese University of Hongkong
    University of Illinois at Chicago
    Hongkong University of Science and Technology (Guangzhou)

    A BSTRACT
    Body starts here.
    """

    affiliations = parser._extract_affiliations(raw_text)

    assert "Tsinghua University" in affiliations
    assert "Beijing University of Posts and Telecommunications" in affiliations
    assert "The Chinese University of Hong Kong" in affiliations
    assert "University of Illinois at Chicago" in affiliations
    assert "Hong Kong University of Science and Technology (Guangzhou)" in affiliations


def test_poster_keypoint_selector_caps_to_ten_by_reading_order(tmp_path, monkeypatch):
    payload = {
        "paper_poster_keypoints": [
            {"id": index, "key_point": f"Poster-worthy claim {index}", "section": "Introduction"}
            for index in range(1, 13)
        ],
        "reading_order": list(range(12, 0, -1)),
    }

    class FakeResponse:
        content = json.dumps(payload)
        input_tokens = 10
        output_tokens = 20

    class FakeAgent:
        def __init__(self, *args, **kwargs):
            pass

        def step(self, message):
            return FakeResponse()

    monkeypatch.setattr("src.agents.poster_keypoint_selector.LangGraphAgent", FakeAgent)
    state = create_state(str(tmp_path / "paper.pdf"))
    state["output_dir"] = str(tmp_path / "output")
    state["raw_text"] = "Full paper text with enough content for keypoint selection."

    result = PosterKeypointSelector()(state)

    assert len(result["paper_poster_keypoints"]) == 10
    assert result["poster_reading_order"] == list(range(1, 11))
    assert result["paper_poster_keypoints"][0]["original_id"] == 12
    assert result["poster_keypoint_selection_report"]["dropped_original_ids"] == [2, 1]
    assert Path(state["output_dir"], "content", "poster_keypoint_selection.json").exists()


def test_poster_keypoint_selector_falls_back_to_structured_sections(tmp_path, monkeypatch):
    class FailingAgent:
        def __init__(self, *args, **kwargs):
            pass

        def step(self, message):
            raise RuntimeError("offline")

    monkeypatch.setattr("src.agents.poster_keypoint_selector.LangGraphAgent", FailingAgent)
    state = create_state(str(tmp_path / "paper.pdf"))
    state["output_dir"] = str(tmp_path / "output")
    state["raw_text"] = "Paper text."
    state["structured_sections"] = {
        "paper_sections": [
            {"section_name": "Introduction", "key_points": ["Problem statement is poster relevant."]},
            {"section_name": "Method", "key_points": ["The framework has a stable arena-based update rule."]},
        ]
    }

    result = PosterKeypointSelector()(state)

    assert len(result["paper_poster_keypoints"]) == 2
    assert result["poster_keypoint_selection_report"]["source"] == "structured_sections_fallback"
    assert result["paper_poster_keypoints"][1]["section"] == "Method"


def test_curator_aligns_story_board_to_keypoints():
    curator = StoryBoardCurator()
    state = create_state("/tmp/paper.pdf")
    state["paper_poster_keypoints"] = [
        {"id": index, "key_point": f"Key contribution or result {index}", "section": "Method" if index <= 6 else "Experiments"}
        for index in range(1, 11)
    ]
    state["poster_reading_order"] = list(range(1, 11))
    story_board = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "method_summary",
                    "section_title": "Long Method Summary Title",
                    "column_assignment": "middle",
                    "vertical_priority": "top",
                    "text_content": ["Existing method detail."],
                    "visual_assets": [],
                }
            ]
        }
    }
    visual_context = {"valid_visual_ids": [], "keypoint_target_count": 10}

    curator._align_sections_to_keypoints(story_board, state, visual_context, classified_visuals={})

    sections = story_board["spatial_content_plan"]["sections"]
    assert len(sections) == 10
    assert [section["keypoint_id"] for section in sections] == list(range(1, 11))
    assert all(section["source_section"] in {"Method", "Experiments"} for section in sections)
    assert all(len(section["section_title"].split()) <= curator.validation_config["max_title_words"] for section in sections)


def test_curator_normalizes_poster_text_items():
    curator = StoryBoardCurator()

    cleaned = curator._clean_poster_text_items(
        [
            "• **Problem:** Dense paper text should become a clean poster item.",
            "    ◦ Nested details should not keep sub-bullet markers.",
            "1) Ordered list markers should be stripped.",
            "",
            "Step 2: Workflow prefixes should be removed.",
            "The results are presented in Table 2.",
        ],
        max_items=6,
    )

    assert cleaned == [
        "**Problem:** Dense paper text should become a clean poster item.",
        "Nested details should not keep sub-bullet markers.",
        "Ordered list markers should be stripped.",
        "Workflow prefixes should be removed.",
    ]
    assert curator._clean_section_title("Main results with table") == "Main Results"
    assert normalize_title_for_poster("Active Geospatial Search For Effcient Tenant Eviction Outreach") == (
        "Active Geospatial Search for Efficient Tenant Eviction Outreach"
    )


def test_curator_groups_keypoints_for_dense_landscape_template(tmp_path):
    curator = StoryBoardCurator()
    capacity_state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_104_landscape",
        width=54,
        height=27,
    )
    capacity_state["output_dir"] = str(tmp_path / "output_dense")
    capacity_state = TemplateCapacityPlanner()(capacity_state)
    state = create_state("/tmp/paper.pdf")
    state["paper_poster_keypoints"] = [
        {"id": 1, "key_point": "Arena comparisons can produce unstable Elo rankings when the same battles are processed in different orders.", "section": "Introduction"},
        {"id": 2, "key_point": "The paper frames stable LLM evaluation as likelihood-based estimation over the complete battle set.", "section": "Introduction"},
        {"id": 3, "key_point": "m-ELO reformulates Elo scoring as maximum likelihood estimation instead of sequential updates.", "section": "Method"},
        {"id": 4, "key_point": "am-ELO extends m-ELO by modeling annotator ability in the pairwise probability function.", "section": "Method"},
        {"id": 5, "key_point": "The framework estimates model scores and annotator abilities jointly from arena records.", "section": "Method"},
        {"id": 6, "key_point": "Experiments compare Elo, m-ELO, and am-ELO on arena-style LLM evaluation data.", "section": "Experiments"},
        {"id": 7, "key_point": "Results show am-ELO achieves lower loss than the baseline Elo-style estimators.", "section": "Results"},
        {"id": 8, "key_point": "Prediction experiments indicate better generalization for am-ELO.", "section": "Results"},
        {"id": 9, "key_point": "Robustness tests analyze perturbations to arena outcomes and annotator behavior.", "section": "Robustness"},
        {"id": 10, "key_point": "The final takeaway is that annotator-aware MLE gives more stable arena-based LLM evaluation.", "section": "Conclusion"},
    ]
    state["poster_reading_order"] = list(range(1, 11))
    story_board = {"spatial_content_plan": {"sections": []}}
    visual_context = {
        "valid_visual_ids": ["figure_2", "figure_3", "table_3"],
        "keypoint_target_count": 10,
        "keypoint_section_target_count": 7,
        "keypoint_grouping_mode": True,
        "requested_layout_template": "cluster_104_landscape",
        "template_fast_mode": True,
        "fast_block_contract": capacity_state["fast_block_contract"],
        "fast_visual_policy": capacity_state["fast_visual_policy"],
        "template_layout": capacity_state["layout_template_metadata"],
        "visual_assets_heights": {
            "figure_2": {"aspect_ratio": 1.7},
            "figure_3": {"aspect_ratio": 2.3},
            "table_3": {"aspect_ratio": 0.94},
        },
    }
    classified_visuals = {
        "key_visual": "figure_2",
        "method_workflow": [],
        "main_results": ["figure_3"],
        "comparative_results": ["table_3"],
        "supporting": [],
    }

    curator._align_sections_to_keypoints(story_board, state, visual_context, classified_visuals)

    sections = story_board["spatial_content_plan"]["sections"]
    source_ids = [keypoint_id for section in sections for keypoint_id in section["source_keypoint_ids"]]
    visual_ids = [
        visual["visual_id"]
        for section in sections
        for visual in section.get("visual_assets", [])
    ]
    assert len(sections) == 7
    assert source_ids == list(range(1, 11))
    assert all(section.get("source_keypoint_ids") for section in sections)
    assert "figure_2" in visual_ids
    assert "figure_3" in visual_ids
    assert "table_3" in visual_ids


def test_curator_groups_keypoints_for_six_slot_landscape_template(tmp_path):
    curator = StoryBoardCurator()
    capacity_state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
    )
    capacity_state["output_dir"] = str(tmp_path / "output_six")
    capacity_state = TemplateCapacityPlanner()(capacity_state)
    state = create_state("/tmp/paper.pdf")
    state["paper_poster_keypoints"] = [
        {"id": index, "key_point": f"am-ELO poster point {index} with stable arena evaluation evidence.", "section": "Method" if index <= 6 else "Results"}
        for index in range(1, 11)
    ]
    state["poster_reading_order"] = list(range(1, 11))
    story_board = {"spatial_content_plan": {"sections": []}}
    visual_context = {
        "valid_visual_ids": ["figure_2", "figure_3", "table_2", "table_3"],
        "keypoint_target_count": 10,
        "keypoint_section_target_count": 6,
        "keypoint_grouping_mode": True,
        "requested_layout_template": "cluster_43_landscape",
        "template_fast_mode": True,
        "fast_block_contract": capacity_state["fast_block_contract"],
        "fast_visual_policy": capacity_state["fast_visual_policy"],
        "template_layout": capacity_state["layout_template_metadata"],
        "visual_assets_heights": {
            "figure_2": {"aspect_ratio": 1.6},
            "figure_3": {"aspect_ratio": 2.0},
            "table_2": {"aspect_ratio": 2.4},
            "table_3": {"aspect_ratio": 1.2},
        },
    }
    classified_visuals = {
        "key_visual": "figure_2",
        "method_workflow": [],
        "main_results": ["figure_3", "table_2"],
        "comparative_results": ["table_3"],
        "supporting": [],
    }

    curator._align_sections_to_keypoints(story_board, state, visual_context, classified_visuals)

    sections = story_board["spatial_content_plan"]["sections"]
    assert len(sections) == 6
    assert [keypoint_id for section in sections for keypoint_id in section["source_keypoint_ids"]] == list(range(1, 11))
    assert [section["preferred_slot_id"] for section in sections] == [
        "slot_1",
        "slot_2",
        "slot_3",
        "slot_4",
        "slot_5",
        "slot_6",
    ]
    visual_ids = [
        visual["visual_id"]
        for section in sections
        for visual in section.get("visual_assets", [])
    ]
    assert {"figure_2", "figure_3", "table_2"}.issubset(set(visual_ids))


def test_template_capacity_planner_builds_landscape_fast_contract(tmp_path):
    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
    )
    state["output_dir"] = str(tmp_path / "output")

    result = TemplateCapacityPlanner()(state)

    contract = result["fast_block_contract"]
    assert result["template_fast_mode"] is True
    assert contract["template_id"] == "cluster_43_landscape"
    assert contract["slot_order"] == ["slot_1", "slot_2", "slot_3", "slot_4", "slot_5", "slot_6"]
    assert len(contract["blocks"]) == 6
    assert contract["by_slot"]["slot_1"]["min_chars"] > 0
    assert contract["by_slot"]["slot_6"]["visual_policy"] == "table_with_callouts"
    assert result["fast_visual_policy"]["figure_count"] == 2
    assert result["fast_visual_policy"]["table_count"] == 2
    assert Path(state["output_dir"], "content", "fast_block_contract.json").exists()


def test_standard_template_preselector_auto_selects_dense_landscape_template(tmp_path):
    state = create_state(str(tmp_path / "paper.pdf"), layout_template="auto", width=54, height=36)
    state["output_dir"] = str(tmp_path / "output")
    state["structured_sections"] = {
        "paper_sections": [
            {"section_name": f"Section {index}", "section_type": "method", "key_points": ["A"]}
            for index in range(6)
        ]
    }
    state["classified_visuals"] = {"main_results": ["table_1"]}
    state["visual_assets"] = {
        "figure_1": {"asset_type": "figure", "aspect": 2.0},
        "figure_2": {"asset_type": "figure", "aspect": 1.5},
        "table_1": {"asset_type": "table", "aspect": 2.4},
    }

    result = StandardTemplatePreselector()(state)

    assert result["resolved_layout_template"] == "cluster_104_landscape"
    assert result["poster_width"] == 54.0
    assert result["poster_height"] == 27.0
    assert result["enable_block_vlm_review"] is True
    assert Path(state["output_dir"], "content", "standard_template_selection_report.json").exists()


def test_template_capacity_planner_builds_standard_template_contract(tmp_path):
    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
    )
    state["output_dir"] = str(tmp_path / "output")

    result = TemplateCapacityPlanner()(state)

    contract = result["fast_block_contract"]
    source_ids = [keypoint_id for block in contract["blocks"] for keypoint_id in block["source_keypoint_ids"]]
    assert result["template_fast_mode"] is True
    assert contract["template_id"] == "cluster_43_landscape"
    assert len(contract["blocks"]) == 6
    assert source_ids == list(range(1, 11))
    assert result["fast_visual_policy"]["figure_slots"]
    assert result["fast_visual_policy"]["table_slots"]


def test_poster_keypoint_selector_prompt_includes_fast_contract():
    selector = PosterKeypointSelector()
    fast_contract = {
        "template_id": "cluster_43_landscape",
        "blocks": [
            {
                "slot_id": "slot_1",
                "slot_role": "Motivation",
                "visual_policy": "text_only",
                "target_chars": 500,
                "source_keypoint_ids": [1, 2],
            }
        ],
    }

    prompt = selector._build_prompt("Paper text.", fast_contract)

    assert "Fast template-first capacity context" in prompt
    assert "cluster_43_landscape" in prompt
    assert "motivation, method/architecture" in prompt


def test_curator_groups_keypoints_for_standard_landscape_template(tmp_path):
    capacity_state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
    )
    capacity_state["output_dir"] = str(tmp_path / "output")
    capacity_state = TemplateCapacityPlanner()(capacity_state)

    curator = StoryBoardCurator()
    state = create_state(str(tmp_path / "paper.pdf"))
    state["paper_poster_keypoints"] = [
        {"id": index, "key_point": f"HAGS poster point {index} with method or result evidence.", "section": "Method" if index <= 6 else "Results"}
        for index in range(1, 11)
    ]
    state["poster_reading_order"] = list(range(1, 11))
    story_board = {"spatial_content_plan": {"sections": []}}
    visual_context = {
        "valid_visual_ids": ["figure_1", "figure_2", "table_1"],
        "keypoint_target_count": 10,
        "keypoint_section_target_count": 6,
        "keypoint_grouping_mode": True,
        "requested_layout_template": "cluster_43_landscape",
        "template_fast_mode": True,
        "fast_block_contract": capacity_state["fast_block_contract"],
        "fast_visual_policy": capacity_state["fast_visual_policy"],
        "template_layout": capacity_state["layout_template_metadata"],
        "visual_assets_heights": {
            "figure_1": {"aspect_ratio": 1.6},
            "figure_2": {"aspect_ratio": 2.0},
            "table_1": {"aspect_ratio": 2.4},
        },
    }
    classified_visuals = {
        "key_visual": "figure_1",
        "method_workflow": [],
        "main_results": ["figure_2", "table_1"],
        "comparative_results": [],
        "supporting": [],
    }

    curator._align_sections_to_keypoints(story_board, state, visual_context, classified_visuals)

    sections = story_board["spatial_content_plan"]["sections"]
    source_ids = [keypoint_id for section in sections for keypoint_id in section["source_keypoint_ids"]]
    visual_ids = [
        visual["visual_id"]
        for section in sections
        for visual in section.get("visual_assets", [])
    ]
    assert len(sections) == 6
    assert source_ids == list(range(1, 11))
    assert all(section.get("preferred_slot_id") for section in sections)
    assert all(section.get("capacity_budget") for section in sections)
    assert {"figure_1", "figure_2", "table_1"}.issubset(set(visual_ids))
    key_holder = next(section for section in sections if section.get("visual_assets") and section["visual_assets"][0]["visual_id"] == "figure_1")
    assert key_holder["column_assignment"] == "middle"
    assert key_holder["vertical_priority"] == "top"


def test_curator_block_template_key_visual_validation_uses_slot_mapping_not_middle_column():
    curator = StoryBoardCurator()
    story_board = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": f"section_{index}",
                    "section_title": f"Block {index}",
                    "column_assignment": "left" if index == 1 else "right",
                    "vertical_priority": "bottom",
                    "source_keypoint_ids": [index],
                    "text_content": ["Evidence-backed point.", "Capacity-aware detail."],
                    "visual_assets": [{"visual_id": "figure_1"}] if index == 1 else [],
                }
                for index in range(1, 5)
            ]
        }
    }
    visual_context = {
        "valid_visual_ids": ["figure_1"],
        "keypoint_target_count": 10,
        "keypoint_section_target_count": 4,
        "keypoint_grouping_mode": True,
        "requested_layout_template": "cluster_43_landscape",
        "visual_assets_heights": {},
    }

    assert curator._validate_story_board(
        story_board,
        classified_visuals={"key_visual": "figure_1"},
        visual_context=visual_context,
    )


def test_curator_portrait_standard_template_keeps_one_key_visual(tmp_path):
    capacity_state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_8_portrait",
        width=36,
        height=50.88,
    )
    capacity_state["output_dir"] = str(tmp_path / "output")
    capacity_state = TemplateCapacityPlanner()(capacity_state)

    curator = StoryBoardCurator()
    state = create_state(str(tmp_path / "paper.pdf"))
    state["paper_poster_keypoints"] = [
        {"id": index, "key_point": f"HAGS portrait keypoint {index}.", "section": "Method" if index <= 6 else "Results"}
        for index in range(1, 11)
    ]
    state["poster_reading_order"] = list(range(1, 11))
    story_board = {"spatial_content_plan": {"sections": []}}
    visual_context = {
        "valid_visual_ids": ["figure_1", "figure_2", "table_1"],
        "keypoint_target_count": 10,
        "keypoint_section_target_count": 4,
        "keypoint_grouping_mode": True,
        "requested_layout_template": "cluster_8_portrait",
        "template_fast_mode": True,
        "fast_block_contract": capacity_state["fast_block_contract"],
        "fast_visual_policy": capacity_state["fast_visual_policy"],
        "template_layout": capacity_state["layout_template_metadata"],
        "visual_assets_heights": {
            "figure_1": {"aspect_ratio": 1.6},
            "figure_2": {"aspect_ratio": 2.0},
            "table_1": {"aspect_ratio": 2.4},
        },
    }
    classified_visuals = {
        "key_visual": "figure_1",
        "method_workflow": [],
        "main_results": ["figure_2", "table_1"],
        "comparative_results": [],
        "supporting": [],
    }

    guidance = curator._template_layout_guidance(visual_context)
    curator._align_sections_to_keypoints(story_board, state, visual_context, classified_visuals)

    sections = story_board["spatial_content_plan"]["sections"]
    visual_ids = [
        visual["visual_id"]
        for section in sections
        for visual in section.get("visual_assets", [])
    ]
    assert "Produce exactly 4 grouped poster sections" in guidance
    assert "Use exactly 1 total visual" in guidance
    assert visual_ids == ["figure_1"]
    assert curator._validate_story_board(story_board, classified_visuals, visual_context)


def test_curator_fast_contract_adds_capacity_budget(tmp_path):
    capacity_state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
    )
    capacity_state["output_dir"] = str(tmp_path / "output")
    capacity_state = TemplateCapacityPlanner()(capacity_state)

    curator = StoryBoardCurator()
    state = create_state(str(tmp_path / "paper.pdf"))
    state["paper_poster_keypoints"] = [
        {"id": index, "key_point": f"Phishing detection keypoint {index} with method or result evidence.", "section": "Method" if index <= 6 else "Results"}
        for index in range(1, 11)
    ]
    state["poster_reading_order"] = list(range(1, 11))
    story_board = {"spatial_content_plan": {"sections": []}}
    visual_context = {
        "valid_visual_ids": ["figure_1", "figure_2", "table_1"],
        "keypoint_target_count": 10,
        "keypoint_section_target_count": 6,
        "keypoint_grouping_mode": True,
        "requested_layout_template": "cluster_43_landscape",
        "template_fast_mode": True,
        "fast_block_contract": capacity_state["fast_block_contract"],
        "fast_visual_policy": capacity_state["fast_visual_policy"],
        "visual_assets_heights": {
            "figure_1": {"aspect_ratio": 1.6},
            "figure_2": {"aspect_ratio": 2.0},
            "table_1": {"aspect_ratio": 2.4},
        },
    }
    classified_visuals = {
        "key_visual": "figure_1",
        "method_workflow": [],
        "main_results": ["figure_2", "table_1"],
        "comparative_results": [],
        "supporting": [],
    }

    curator._align_sections_to_keypoints(story_board, state, visual_context, classified_visuals)

    sections = story_board["spatial_content_plan"]["sections"]
    by_slot = {section["preferred_slot_id"]: section for section in sections}
    assert len(sections) == 6
    assert by_slot["slot_1"]["target_chars"] > 0
    assert by_slot["slot_5"]["min_chars"] > 0
    assert by_slot["slot_6"]["capacity_budget"]["visual_policy"] == "table_with_callouts"
    assert [keypoint_id for section in sections for keypoint_id in section["source_keypoint_ids"]] == list(range(1, 11))


def test_template_planner_uses_slot_mapping_and_preserves_method_visual():
    planner = TemplateBlockPlanner()
    cfg = load_config()
    template = LayoutTemplates(
        54,
        27,
        margin=cfg["layout"]["poster_margin"],
        col_gap=cfg["layout"]["column_spacing"],
    ).get_template("cluster_43_landscape")
    regions = template["regions"]
    sections = [
        {"section_id": "why", "section_title": "Why ELO Fails", "text_content": ["Motivation"], "visual_assets": [], "content_role": "overview", "preferred_slot_id": "slot_1"},
        {"section_id": "stable", "section_title": "Stable Estimation", "text_content": ["m-ELO"], "visual_assets": [{"visual_id": "figure_3"}], "content_role": "method", "preferred_slot_id": "slot_2"},
        {"section_id": "arena", "section_title": "Arena Setting", "text_content": ["Arena diagram"], "visual_assets": [{"visual_id": "figure_2"}], "content_role": "method", "preferred_slot_id": "slot_3"},
        {"section_id": "annotator", "section_title": "Annotator-Aware ELO", "text_content": ["Annotators"], "visual_assets": [], "content_role": "method", "preferred_slot_id": "slot_6"},
        {"section_id": "robust", "section_title": "Robustness Tests", "text_content": ["Robustness"], "visual_assets": [{"visual_id": "table_3"}], "content_role": "results", "preferred_slot_id": "slot_5"},
        {"section_id": "main", "section_title": "Main Results", "text_content": ["Results"], "visual_assets": [], "content_role": "results", "preferred_slot_id": "slot_4"},
    ]

    assigned = planner._assign_sections_to_regions(
        sections,
        regions,
        hero_section=sections[0],
        hero_region_id="slot_1",
        preserve_order=True,
        template_id="cluster_43_landscape",
    )

    by_id = {section["section_id"]: section for section in assigned}
    assert by_id["main"]["slot_id"] == "slot_4"
    assert by_id["robust"]["slot_id"] == "slot_5"
    assert by_id["annotator"]["slot_id"] == "slot_6"
    assert by_id["arena"]["slot_id"] == "slot_3"
    assert by_id["arena"]["visual_assets"][0]["visual_id"] == "figure_2"


def test_template_planner_preserves_fast_policy_table_in_low_density_slot():
    planner = TemplateBlockPlanner()
    cfg = load_config()
    template = LayoutTemplates(
        54,
        27,
        margin=cfg["layout"]["poster_margin"],
        col_gap=cfg["layout"]["column_spacing"],
    ).get_template("cluster_43_landscape")
    regions = template["regions"]
    sections = [
        {"section_id": "motivation", "section_title": "Why AGS", "text_content": ["Motivation"], "visual_assets": [], "content_role": "foundation", "preferred_slot_id": "slot_1"},
        {"section_id": "method", "section_title": "Search Tension", "text_content": ["Method"], "visual_assets": [{"visual_id": "figure_2"}], "content_role": "method", "preferred_slot_id": "slot_2"},
        {"section_id": "flow", "section_title": "HAGS Flow", "text_content": ["Flow"], "visual_assets": [{"visual_id": "figure_1"}], "content_role": "method", "preferred_slot_id": "slot_3"},
        {"section_id": "learns", "section_title": "How It Learns", "text_content": ["Learning"], "visual_assets": [], "content_role": "method", "preferred_slot_id": "slot_4"},
        {"section_id": "setup", "section_title": "Data Setup", "text_content": ["Setup"], "visual_assets": [], "content_role": "results", "preferred_slot_id": "slot_5"},
        {
            "section_id": "results",
            "section_title": "Key Results",
            "text_content": ["Results"],
            "visual_assets": [{"visual_id": "table_2"}],
            "content_role": "results",
            "preferred_slot_id": "slot_6",
            "capacity_budget": {"visual_policy": "table_with_callouts"},
        },
    ]

    assigned = planner._assign_sections_to_regions(
        sections,
        regions,
        hero_section=sections[1],
        hero_region_id="slot_2",
        preserve_order=True,
        template_id="cluster_43_landscape",
    )

    by_id = {section["section_id"]: section for section in assigned}
    assert by_id["results"]["slot_id"] == "slot_6"
    assert by_id["results"]["visual_assets"][0]["visual_id"] == "table_2"


def test_affiliation_logo_agent_creates_placeholder_when_download_fails(tmp_path, monkeypatch):
    state = create_state(str(tmp_path / "paper.pdf"), enable_affiliation_logos=True)
    state["output_dir"] = str(tmp_path / "output")
    state["affiliations"] = ["Example Research University"]

    agent = AffiliationLogoAgent()
    agent.config["include_placeholders"] = True
    monkeypatch.setattr(agent, "_download_clearbit_logo", lambda domain, output_path: None)
    monkeypatch.setattr(agent, "_download_wikidata_logo", lambda institution, output_path: None)
    monkeypatch.setattr(agent, "_download_known_commons_logo", lambda institution, output_path: None)

    result = agent(state)

    logos = result["affiliation_logos"]
    assert len(logos) == 1
    assert logos[0]["status"] == "placeholder"
    assert Path(logos[0]["logo_path"]).exists()
    assert (Path(state["output_dir"]) / "content" / "affiliation_logos.json").exists()


def test_layout_agent_places_affiliation_logos_in_title_right_region(tmp_path):
    logo_paths = []
    for index in range(3):
        path = tmp_path / f"logo_{index}.png"
        Image.new("RGBA", (300, 120), (255, 255, 255, 255)).save(path)
        logo_paths.append(str(path))

    state = create_state("/tmp/paper.pdf", enable_affiliation_logos=True)
    state["affiliation_logos"] = [
        {
            "institution": f"Institution {index}",
            "logo_path": path,
            "domain": None,
            "source": "test",
            "aspect": 2.5,
        }
        for index, path in enumerate(logo_paths)
    ]

    agent = LayoutAgent()
    template = agent._resolve_template_layout(state)
    elements = agent._create_logo_elements(state, state["poster_width"], template)

    assert len(elements) == 3
    assert {element["type"] for element in elements} == {"institution_logo"}
    title = agent._create_title_element(state, state["poster_width"], template["header"]["h"], template)
    assert all(element["x"] >= title["x"] + title["width"] for element in elements)


def test_layout_agent_avoids_title_conference_logo_overlap_for_portrait_templates(tmp_path):
    logo_path = tmp_path / "conference.png"
    Image.new("RGBA", (900, 420), (20, 80, 160, 255)).save(logo_path)

    agent = LayoutAgent()
    for template_id in list_block_template_ids():
        info = get_block_template_info(template_id)
        canvas = info["recommended_canvas_size"]
        state = create_state("/tmp/paper.pdf", layout_template=template_id)
        state["poster_width"] = canvas["width"]
        state["poster_height"] = canvas["height"]
        state["resolved_layout_template"] = template_id
        state["template_layout_mode"] = "template_prior"
        state["logo_path"] = str(logo_path)

        template = agent._resolve_template_layout(state)
        title = agent._create_title_element(state, state["poster_width"], template["header"]["h"], template)
        logo = next(
            element for element in agent._create_logo_elements(state, state["poster_width"], template)
            if element["type"] == "conf_logo"
        )

        title_right = title["x"] + title["width"]
        logo_left = logo["x"]
        assert title_right <= logo_left


def test_layout_agent_new_landscape_header_keeps_title_and_logo_zone(tmp_path):
    conf_path = tmp_path / "conference.png"
    aff_path = tmp_path / "affiliation.png"
    Image.new("RGBA", (900, 420), (20, 80, 160, 255)).save(conf_path)
    Image.new("RGBA", (700, 700), (160, 40, 60, 255)).save(aff_path)

    state = create_state(
        "/tmp/paper.pdf",
        layout_template="cluster_43_landscape",
        width=54,
        height=27,
        logo_path=str(conf_path),
        aff_logo_path=str(aff_path),
    )
    state["resolved_layout_template"] = "cluster_43_landscape"
    state["template_layout_mode"] = "template_prior"

    agent = LayoutAgent()
    template = agent._resolve_template_layout(state)
    title_box, logo_box = agent._header_title_logo_boxes(state, template)
    title = agent._create_title_element(state, state["poster_width"], template["header"]["h"], template)

    assert title_box["w"] >= template["header"]["w"] * 0.60
    assert logo_box["w"] >= template["header"]["w"] * 0.28
    assert title_box["x"] + title_box["w"] <= logo_box["x"]
    assert title["font_size"] >= 96
    assert title["author_font_size"] <= 72


def test_layout_agent_section_title_uses_navy_band_wordart():
    agent = LayoutAgent()
    state = create_state("/tmp/paper.pdf")
    section = {
        "section_id": "method",
        "section_title": "Method",
        "column_assignment": "slot_3",
        "slot_id": "slot_3",
    }

    elements = agent._create_section_title_design(section, column_x=1.0, start_y=2.0, column_width=6.0, state=state)

    bar = next(element for element in elements if element["type"] == "title_accent_block")
    title = next(element for element in elements if element["type"] == "section_title")
    assert bar["x"] == pytest.approx(1.0)
    assert bar["width"] == pytest.approx(6.0)
    assert bar["color"] == "#06134A"
    assert title["section_title"] == "3. Method"
    assert title["font_family"] == "Georgia"
    assert title["font_color"] == "#FFFFFF"
    assert title["wordart_style"]["name"] == "navy_band_serif"


def test_section_title_designer_emits_navy_band_template():
    state = create_state("/tmp/paper.pdf")
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {"section_id": "method", "section_title": "Method"},
            ]
        }
    }
    state["color_scheme"] = {"theme": "#335f91", "mono_light": "#8AA0BA", "mono_dark": "#001E44"}

    result = SectionTitleDesigner()(state)

    design = result["section_title_design"]["section_title_design"]
    assert design["selected_template"] == "navy_band_wordart"
    application = design["section_applications"][0]
    assert application["accent_styling"]["type"] == "full_width_bar"
    assert application["title_styling"]["font_family"] == "Georgia"


def test_layout_agent_uses_manual_aff_logo_with_conference_logo(tmp_path):
    conf_path = tmp_path / "conference.png"
    aff_path = tmp_path / "affiliation.png"
    Image.new("RGBA", (900, 463), (20, 80, 160, 255)).save(conf_path)
    Image.new("RGBA", (1200, 1200), (160, 40, 60, 255)).save(aff_path)

    state = create_state("/tmp/paper.pdf", layout_template="cluster_3_portrait", width=36, height=50.88)
    state["logo_path"] = str(conf_path)
    state["aff_logo_path"] = str(aff_path)

    agent = LayoutAgent()
    template = agent._resolve_template_layout(state)
    elements = agent._create_logo_elements(state, state["poster_width"], template)

    assert {element["type"] for element in elements} == {
        "conf_logo",
        "institution_logo",
        "logo_divider",
    }
    institution_logo = next(element for element in elements if element["type"] == "institution_logo")
    conference_logo = next(element for element in elements if element["type"] == "conf_logo")
    assert institution_logo["image_path"] == str(aff_path)
    assert institution_logo["source"] == "manual"
    assert institution_logo["height"] >= 1.8
    assert conference_logo["height"] >= 2.4


def test_layout_agent_supports_five_affiliation_logos_without_overflow(tmp_path):
    logo_paths = []
    for index in range(5):
        path = tmp_path / f"logo_{index}.png"
        Image.new("RGBA", (420, 120), (20, 50, 120, 255)).save(path)
        logo_paths.append(str(path))

    state = create_state("/tmp/paper.pdf", enable_affiliation_logos=True)
    state["affiliation_logos"] = [
        {
            "institution": f"Institution {index}",
            "logo_path": path,
            "domain": None,
            "source": "test",
            "aspect": 3.5,
        }
        for index, path in enumerate(logo_paths)
    ]

    agent = LayoutAgent()
    template = agent._resolve_template_layout(state)
    elements = agent._create_logo_elements(state, state["poster_width"], template)
    logos = [element for element in elements if element.get("type") == "institution_logo"]
    region = agent._title_logo_region(template, False)

    assert len(logos) == 5
    assert all(region["x"] <= logo["x"] for logo in logos)
    assert all(logo["x"] + logo["width"] <= region["x"] + region["w"] + 1e-6 for logo in logos)
    assert all(region["y"] <= logo["y"] for logo in logos)
    assert all(logo["y"] + logo["height"] <= region["y"] + region["h"] + 1e-6 for logo in logos)


def test_micro_layout_refiner_keeps_logo_divider_as_global_element():
    state = create_state("/tmp/paper.pdf", layout_template="three_column_postergen")
    state["styled_layout"] = [
        {
            "type": "title",
            "x": 1.0,
            "y": 1.0,
            "width": 32.0,
            "height": 4.0,
            "priority": 1.0,
        },
        {
            "type": "logo_divider",
            "x": 42.0,
            "y": 1.5,
            "width": 0.04,
            "height": 3.5,
            "priority": 0.85,
        },
        {
            "type": "section_container",
            "section_id": "left::s1",
            "lane_id": "left",
            "x": 1.0,
            "y": 7.0,
            "width": 16.0,
            "height": 4.0,
            "priority": 0.1,
        },
        {
            "type": "text",
            "id": "left::s1_text",
            "section_id": "left::s1",
            "x": 1.3,
            "y": 7.6,
            "width": 15.4,
            "height": 2.5,
            "content": "Short text",
            "font_size": 40,
            "priority": 0.5,
        },
    ]

    result = MicroLayoutRefiner()(state)
    divider = next(element for element in result["styled_layout"] if element.get("type") == "logo_divider")

    assert divider["y"] == 1.5


def test_renderer_uses_element_path_for_institution_logo(tmp_path):
    logo_path = tmp_path / "logo.png"
    Image.new("RGBA", (100, 80), (255, 255, 255, 255)).save(logo_path)

    calls = []

    class FakeDirector:
        def add_image(self, *args, **kwargs):
            calls.append((args, kwargs))

    renderer = Renderer()
    renderer.director = FakeDirector()
    element = {
        "type": "institution_logo",
        "x": 1,
        "y": 1,
        "width": 2,
        "height": 1,
        "image_path": str(logo_path),
    }

    renderer._render_institution_logo(None, element, create_state("/tmp/paper.pdf"))

    assert calls
    assert calls[0][0][0] == str(logo_path)


def test_renderer_trims_logo_whitespace(tmp_path):
    logo_path = tmp_path / "logo_with_margin.png"
    image = Image.new("RGBA", (400, 400), (255, 255, 255, 255))
    for x in range(150, 250):
        for y in range(120, 280):
            image.putpixel((x, y), (20, 60, 120, 255))
    image.save(logo_path)

    renderer = Renderer()
    renderer._render_output_dir = tmp_path / "output"
    trimmed_path = Path(renderer._trim_logo_whitespace(str(logo_path)))

    assert trimmed_path.exists()
    with Image.open(trimmed_path) as trimmed:
        assert trimmed.size[0] < 160
        assert trimmed.size[1] < 220


def test_renderer_draws_explicit_section_container_fill_for_standard_templates():
    calls = []

    class FakeDirector:
        def add_shape(self, *args, **kwargs):
            calls.append((args, kwargs))

    renderer = Renderer()
    renderer.director = FakeDirector()
    element = {
        "type": "section_container",
        "x": 1,
        "y": 2,
        "width": 5,
        "height": 4,
        "fill_color": "#F0F6FF",
        "border_color": "#C9DDF5",
        "border_width": 0.9,
        "border_style": "dashed",
        "shadow": {"enabled": True, "color": "#000000", "alpha": 0.16},
    }

    renderer._render_section_container(None, element, create_state("/tmp/paper.pdf"))

    assert calls
    assert calls[0][1]["fill_color"] == "#F0F6FF"
    assert calls[0][1]["border_color"] == "#C9DDF5"
    assert calls[0][1]["border_style"] == "dashed"
    assert calls[0][1]["shadow"] == element["shadow"]


def test_pptx_director_add_shape_applies_outer_shadow_xml():
    director = PPTXDirector()

    shape = director.add_shape(
        MSO_SHAPE.RECTANGLE,
        0.1,
        0.1,
        1.0,
        1.0,
        fill_color="#F1F2F4",
        shadow={"enabled": True, "color": "#000000", "alpha": 0.16, "blur_pt": 5, "distance_pt": 2.4},
    )

    xml = shape._element.xml
    assert "outerShdw" in xml
    assert 'alpha val="16000"' in xml


def test_renderer_draws_background_image_before_layout(tmp_path):
    background_path = tmp_path / "background.png"
    Image.new("RGB", (100, 140), (245, 247, 250)).save(background_path)
    calls = []

    class FakeDirector:
        def set_slide_dimensions(self, *args, **kwargs):
            pass

        @property
        def slide(self):
            return None

        def add_image(self, *args, **kwargs):
            calls.append((args, kwargs))

        def save(self, *args, **kwargs):
            pass

    state = create_state("/tmp/paper.pdf", width=36, height=50.88)
    state["background_image_path"] = str(background_path)
    state["styled_layout"] = [
        {
            "type": "title",
            "x": 1,
            "y": 1,
            "width": 10,
            "height": 2,
            "content": "Title",
        }
    ]

    renderer = Renderer()
    renderer.director = FakeDirector()
    renderer._render_background_image(None, state)

    assert calls[0][0][:5] == (str(background_path), 0, 0, 36, 50.88)
    assert calls[0][1]["keep_aspect_ratio"] is False


def test_renderer_separates_title_and_authors_with_physical_gap():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    renderer = Renderer()
    renderer.styling_interfaces = {"font_sizes": {"title": 100, "authors": 72}}
    before_shape_count = len(slide.shapes)
    element = {
        "type": "title",
        "x": 1,
        "y": 1,
        "width": 12,
        "height": 4,
        "content": "A Long Research Poster Title\nA. Author, B. Author",
        "font_size": 100,
        "author_font_size": 72,
    }

    renderer._render_title(slide, element, create_state("/tmp/paper.pdf"))

    assert len(slide.shapes) == before_shape_count + 3
    title_box = slide.shapes[-2]
    author_box = slide.shapes[-1]
    actual_gap = author_box.top.inches - (title_box.top.inches + title_box.height.inches)
    assert actual_gap == pytest.approx(16 / 72, abs=0.01)


def test_block_template_registry_exposes_cluster_templates():
    template_ids = set(list_block_template_ids())

    assert {"cluster_2_landscape", "cluster_43_landscape", "cluster_3_portrait", "cluster_8_portrait"}.issubset(template_ids)
    assert "cluster_27_landscape" in template_ids
    assert "cluster_27_portrait" in template_ids


def test_block_template_layout_identifies_header_and_content_slots():
    layout = load_block_template_layout("cluster_8_portrait", 36, 51, margin=1.0)

    assert layout["layout_mode"] == "template_prior"
    assert layout["orientation"] == "portrait"
    assert layout["template_aspect_ratio"] < 1
    assert layout["header_slot"]["slot_id"] == "slot_0"
    assert len(layout["content_slots"]) == 4
    assert len(layout["lanes"]) == 4
    assert all(slot["x"] + slot["w"] <= 36.05 for slot in layout["content_slots"])
    assert all(slot["y"] + slot["h"] <= 51.05 for slot in layout["content_slots"])
    assert all(slot["slot_id"] != layout["header_slot"]["slot_id"] for slot in layout["content_slots"])


def test_new_landscape_template_softly_expands_into_gutters():
    layout = load_block_template_layout("cluster_43_landscape", 54, 27, margin=1.0)
    slot_by_id = {slot["slot_id"]: slot for slot in layout["content_slots"]}
    report = layout["gap_absorption_report"]

    assert report["enabled"] is True
    assert report["mode"] == "equal_split"
    assert len(report["absorptions"]) >= 2
    assert len(report["edge_expansions"]) >= 2
    assert report["total_area_gain"] > 0
    assert any(absorption["orientation"] == "vertical" for absorption in report["absorptions"])
    assert any(absorption["orientation"] == "horizontal" for absorption in report["absorptions"])
    assert any(slot.get("gap_absorbed") is True for slot in slot_by_id.values())
    assert slot_by_id["slot_3"]["x"] + slot_by_id["slot_3"]["w"] <= 53.01


def test_standard_template_soft_geometry_is_enabled():
    layout = load_block_template_layout("cluster_43_landscape", 54, 27, margin=1.0)
    report = layout["gap_absorption_report"]

    assert report["enabled"] is True
    assert report["mode"] == "equal_split"
    assert report["area_after"] > report["area_before"]
    assert report["total_area_gain"] > 0
    assert report["edge_expansions"]
    assert {
        "affected_slot_ids",
        "orientation",
        "original_gap_inches",
        "split_boundary",
        "left_or_upper_gain",
        "right_or_lower_gain",
    }.issubset(report["absorptions"][0])


def test_equal_split_gap_absorption_splits_horizontal_neighbors_at_midline():
    raw = {
        "aspect_ratio": 2.0,
        "slots": [
            {"slot_id": 0, "bbox": [0, 0, 1000, 120]},
            {"slot_id": 1, "bbox": [0, 200, 400, 700]},
            {"slot_id": 2, "bbox": [600, 200, 1000, 700]},
        ],
    }

    layout = build_runtime_template(raw, "cluster_43_landscape", 20, 10, margin=1.0)
    slot_1, slot_2 = layout["content_slots"]
    report = layout["gap_absorption_report"]

    assert report["enabled"] is True
    assert report["mode"] == "equal_split"
    assert len(report["absorptions"]) == 1
    assert report["absorptions"][0]["orientation"] == "vertical"
    assert slot_1["x"] + slot_1["w"] == pytest.approx(slot_2["x"], abs=0.001)
    assert slot_1["w"] > 7.2
    assert slot_2["w"] > 7.2
    assert slot_1["slot_id"] == "slot_1"
    assert slot_2["slot_id"] == "slot_2"


def test_equal_split_gap_absorption_expands_outer_edges_to_safe_bounds():
    layout = load_block_template_layout("cluster_43_landscape", 54, 27, margin=1.0)
    slot_by_id = {slot["slot_id"]: slot for slot in layout["content_slots"]}
    report = layout["gap_absorption_report"]

    assert any(item["edge"] == "bottom" for item in report["edge_expansions"])
    assert any(item["edge"] == "left" for item in report["edge_expansions"])
    assert any(item["edge"] == "right" for item in report["edge_expansions"])
    assert slot_by_id["slot_4"]["x"] == pytest.approx(1.0)
    assert slot_by_id["slot_4"]["y"] + slot_by_id["slot_4"]["h"] == pytest.approx(26.0)
    assert slot_by_id["slot_6"]["x"] + slot_by_id["slot_6"]["w"] == pytest.approx(53.0)
    assert slot_by_id["slot_6"]["y"] + slot_by_id["slot_6"]["h"] == pytest.approx(26.0)


def test_equal_split_gap_absorption_splits_vertical_neighbors_at_midline():
    raw = {
        "aspect_ratio": 1.0,
        "slots": [
            {"slot_id": 0, "bbox": [0, 0, 1000, 120]},
            {"slot_id": 1, "bbox": [0, 200, 1000, 480]},
            {"slot_id": 2, "bbox": [0, 680, 1000, 1000]},
        ],
    }

    layout = build_runtime_template(raw, "cluster_43_landscape", 20, 10, margin=1.0)
    slot_1, slot_2 = layout["content_slots"]
    report = layout["gap_absorption_report"]

    assert report["enabled"] is True
    assert report["mode"] == "equal_split"
    assert len(report["absorptions"]) == 1
    assert report["absorptions"][0]["orientation"] == "horizontal"
    assert slot_1["y"] + slot_1["h"] == pytest.approx(slot_2["y"], abs=0.001)
    assert slot_1["h"] > 2.24
    assert slot_2["h"] > 2.56
    assert slot_1["slot_id"] == "slot_1"
    assert slot_2["slot_id"] == "slot_2"


def test_layout_agent_selects_single_primary_block_background():
    state = create_state("/tmp/paper.pdf", layout_template="cluster_3_portrait", width=36, height=50.88)
    state["color_scheme"] = {"theme": "#0057B8", "contrast": "#7F4B13"}
    agent = LayoutAgent()
    agent.config["selective_block_backgrounds"]["enabled"] = True
    agent.config["selective_block_backgrounds"]["max_highlight_blocks"] = 1
    template = agent._resolve_template_layout(state)
    column_assignments = [
        {
            "column_name": "slot_1",
            "sections": [
                {
                    "section_id": "stable_framework",
                    "section_title": "Stable Framework",
                    "content_role": "method",
                    "visual_assets": [{"visual_id": "figure_2"}],
                }
            ],
        },
        {
            "column_name": "slot_2",
            "sections": [
                {
                    "section_id": "main_results",
                    "section_title": "Main Results",
                    "content_role": "results",
                    "visual_assets": [{"visual_id": "figure_3"}],
                }
            ],
        },
        {
            "column_name": "slot_3",
            "sections": [
                {
                    "section_id": "background",
                    "section_title": "Background",
                    "content_role": "overview",
                    "visual_assets": [],
                }
            ],
        },
    ]

    selected = agent._select_highlight_section_ids(column_assignments, state, template)

    assert selected == {"main_results": 0}


def test_layout_agent_applies_primary_gray_highlight_panel():
    state = create_state("/tmp/paper.pdf", layout_template="cluster_3_portrait", width=36, height=50.88)
    state["color_scheme"] = {"theme": "#0057B8", "contrast": "#7F4B13"}
    agent = LayoutAgent()
    agent.config["selective_block_backgrounds"]["enabled"] = True
    container = {"section_id": "stable_framework", "priority": 0.1}
    section = {"section_id": "stable_framework", "content_role": "method"}

    agent._apply_selective_highlight_panel(container, section, state, {"stable_framework": 0})

    assert container["highlight_panel"] is True
    assert container["fill_color"] == "#EEF0F2"
    assert "border_color" not in container
    assert "border_style" not in container
    assert container["priority"] <= 0.08


def test_layout_agent_frames_normal_and_support_blocks():
    state = create_state("/tmp/paper.pdf", layout_template="cluster_3_portrait", width=36, height=50.88)
    agent = LayoutAgent()
    agent.config["selective_block_backgrounds"]["enabled"] = True
    agent.config["selective_block_backgrounds"]["frame_all_blocks"] = True
    normal = {
        "section_id": "method_details",
        "template_prior": True,
        "priority": 0.1,
    }
    support = {
        "section_id": "robustness_checks",
        "template_prior": True,
        "priority": 0.1,
    }

    agent._apply_selective_block_frame_style(
        normal,
        {"section_id": "method_details", "section_title": "Method Details", "content_role": "method"},
        state,
        {},
    )
    agent._apply_selective_block_frame_style(
        support,
        {"section_id": "robustness_checks", "section_title": "Robustness", "content_role": "overview"},
        state,
        {},
    )

    assert "fill_color" not in normal
    assert normal["border_color"] == "#D2D6DC"
    assert normal["border_style"] == "solid"
    assert support["border_color"] == "#D8DDE3"
    assert support["border_style"] == "dashed"


def test_create_state_uses_draft_stage_when_post_render_pass_is_enabled():
    assert create_state("/tmp/paper.pdf")["render_stage"] == "final"
    assert create_state("/tmp/paper.pdf", enable_generated_background=True)["render_stage"] == "draft"
    assert create_state("/tmp/paper.pdf", enable_vlm_layout_review=True)["render_stage"] == "draft"
    assert create_state("/tmp/paper.pdf", enable_visual_legibility_review=True)["render_stage"] == "draft"
    assert create_state("/tmp/paper.pdf", enable_block_vlm_review=True)["render_stage"] == "draft"


def test_background_image_agent_prompt_is_background_only():
    state = create_state("/tmp/paper.pdf")
    state["color_scheme"] = {"theme": "#0057B8", "mono_light": "#E6EAEF"}
    state["background_palette"] = "light_blue"
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "main_results",
                    "section_title": "Main Results",
                    "content_role": "results",
                }
            ]
        }
    }
    prompt = BackgroundImageAgent()._build_prompt(state)

    assert "BACKGROUND ONLY" in prompt
    assert "no text" in prompt
    assert "not plain white" in prompt
    assert "Selected background palette: light_blue" in prompt
    assert "#0057B8" in prompt


def test_background_image_agent_switches_palette_prompt():
    state = create_state("/tmp/paper.pdf", background_palette="light_gray")

    prompt = BackgroundImageAgent()._build_prompt(state)

    assert "pale neutral gray" in prompt
    assert "Selected background palette: light_gray" in prompt


def test_background_image_agent_placeholder_fallback_keeps_pipeline_success(tmp_path, monkeypatch):
    def fake_generate_image(self, prompt, width, height, output_path):
        Image.new("RGB", (width, height), color=(200, 200, 200)).save(output_path)
        return output_path

    monkeypatch.setattr("src.agents.background_image_agent.ImageTools.generate_image", fake_generate_image)
    state = create_state(str(tmp_path / "paper.pdf"), enable_generated_background=True, background_palette="light_blue")
    state["output_dir"] = str(tmp_path / "output")
    state["color_scheme"] = {"theme": "#0057B8", "mono_light": "#E6EAEF"}
    agent = BackgroundImageAgent()
    agent.background_config["width_px"] = 160
    agent.background_config["height_px"] = 120

    result = agent(state)

    assert result["errors"] == []
    assert Path(result["background_image_path"]).exists()
    assert result["background_image_report"]["palette"] == "light_blue"
    assert result["background_image_report"]["used_procedural_fallback"] is True


def test_background_image_agent_uses_poster_preview_as_reference(tmp_path, monkeypatch):
    def fake_edit_image(self, image_path, prompt, output_path):
        Image.new("RGB", (160, 120), color=(242, 248, 255)).save(output_path)
        assert image_path.endswith("draft.png")
        assert "provided poster image" in prompt
        return output_path

    monkeypatch.setattr("src.agents.background_image_agent.ImageTools.edit_image", fake_edit_image)
    state = create_state(str(tmp_path / "paper.pdf"), enable_generated_background=True, background_palette="light_blue")
    state["output_dir"] = str(tmp_path / "output")
    state["color_scheme"] = {"theme": "#0057B8", "mono_light": "#E6EAEF"}
    preview_path = tmp_path / "draft.png"
    Image.new("RGB", (160, 120), color=(255, 255, 255)).save(preview_path)
    state["poster_preview_path"] = str(preview_path)
    agent = BackgroundImageAgent()
    agent.background_config["width_px"] = 160
    agent.background_config["height_px"] = 120
    agent.background_config["procedural_only"] = False

    result = agent(state)

    assert result["errors"] == []
    assert Path(result["background_image_path"]).exists()
    assert result["background_image_report"]["generation_mode"] == "poster_conditioned_image_api_with_procedural_fallback"
    assert result["background_image_report"]["reference_poster_path"] == str(preview_path)


def test_background_image_agent_procedural_only_skips_image_api(tmp_path, monkeypatch):
    def fail_generate_image(self, prompt, width, height, output_path):
        raise AssertionError("image generation API should not be called")

    monkeypatch.setattr("src.agents.background_image_agent.ImageTools.generate_image", fail_generate_image)
    state = create_state(str(tmp_path / "paper.pdf"), enable_generated_background=True, background_palette="light_blue")
    state["output_dir"] = str(tmp_path / "output")
    agent = BackgroundImageAgent()
    agent.background_config["procedural_only"] = True
    agent.background_config["width_px"] = 160
    agent.background_config["height_px"] = 120

    result = agent(state)

    assert result["errors"] == []
    assert Path(result["background_image_path"]).exists()
    assert result["background_image_report"]["generation_mode"] == "procedural_only"
    assert result["background_image_report"]["raw_path"] == ""


def test_background_image_agent_matches_landscape_poster_aspect():
    state = create_state("/tmp/paper.pdf", width=54, height=27, enable_generated_background=True)
    agent = BackgroundImageAgent()
    agent.background_config["width_px"] = 160
    agent.background_config["height_px"] = 120

    assert agent._background_dimensions(state) == (160, 80)


def test_image_tools_failover_retries_each_base_url(tmp_path, monkeypatch):
    import base64
    import io

    img = Image.new("RGB", (1, 1), color=(240, 248, 255))
    buffer = io.BytesIO()
    img.save(buffer, format="PNG")
    b64 = base64.b64encode(buffer.getvalue()).decode("ascii")
    calls = []

    class FakeResponse:
        def raise_for_status(self):
            return None

        def json(self):
            return {"data": [{"b64_json": b64}]}

    def fake_post(url, **kwargs):
        calls.append(url)
        if url.startswith("https://first.example"):
            raise RuntimeError("temporary upstream failure")
        return FakeResponse()

    monkeypatch.setattr("src.tools.image_api.requests.post", fake_post)
    tool = ImageTools(
        api_key="test-key",
        base_url="https://first.example/v1, https://second.example/v1",
        model="gpt-image-2",
        retry_attempts=2,
        retry_delay=0,
    )
    output_path = tmp_path / "generated.png"

    result = tool.generate_image("plain academic background", width=1600, height=900, output_path=str(output_path))

    assert result == str(output_path)
    assert output_path.exists()
    assert calls == [
        "https://first.example/v1/images/generations",
        "https://first.example/v1/images/generations",
        "https://second.example/v1/images/generations",
    ]


def test_image_tools_gpt_image_uses_supported_aspect_size():
    tool = ImageTools(
        api_key="test-key",
        base_url="https://example.test/v1",
        model="gpt-image-2",
        retry_attempts=1,
        retry_delay=0,
    )

    assert tool._request_size(2035, 1018) == "1536x1024"
    assert tool._request_size(1018, 2035) == "1024x1536"
    assert tool._request_size(1024, 1024) == "1024x1024"


def test_image_tools_base_url_list_takes_priority_over_legacy_env(monkeypatch):
    monkeypatch.setenv("IMAGE_BASE_URLS", "https://first.example/v1, https://second.example/v1")
    monkeypatch.setenv("IMAGE_BASE_URL", "https://legacy-image.example/v1")
    monkeypatch.setenv("VLM_BASE_URL", "https://legacy-vlm.example/v1")

    tool = ImageTools(api_key="test-key", model="gpt-image-2", retry_attempts=1, retry_delay=0)

    assert tool.base_urls == ["https://first.example/v1", "https://second.example/v1"]


def test_template_block_planner_matches_block_count_to_content_slots(monkeypatch):
    json_response = """{
      "blocks": [
        {"target_title": "Overview", "target_bullets": ["Problem framing.", "Core idea."]},
        {"target_title": "Method", "target_bullets": ["Framework details.", "Optimization flow."]},
        {"target_title": "Results", "target_bullets": ["Main result.", "Comparison summary."]},
        {"target_title": "Takeaways", "target_bullets": ["Deployment note.", "Conclusion."]}
      ]
    }"""

    class FakeResponse:
        input_tokens = 1
        output_tokens = 1
        content = json_response

    class FakeAgent:
        def __init__(self, *args, **kwargs):
            pass

        def step(self, message):
            return FakeResponse()

    monkeypatch.setattr("src.agents.template_block_planner.LangGraphAgent", FakeAgent)

    state = create_state("/tmp/paper.pdf", layout_template="cluster_2_landscape", width=54, height=27)
    state["resolved_layout_template"] = "cluster_2_landscape"
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "overview",
                    "section_title": "Overview",
                    "column_assignment": "left",
                    "vertical_priority": "top",
                    "text_content": ["Problem framing.", "Core idea."],
                    "visual_assets": [],
                },
                {
                    "section_id": "method",
                    "section_title": "Method",
                    "column_assignment": "middle",
                    "vertical_priority": "top",
                    "text_content": ["Framework details.", "Optimization flow."],
                    "visual_assets": [{"visual_id": "figure_1"}],
                },
                {
                    "section_id": "results",
                    "section_title": "Results",
                    "column_assignment": "right",
                    "vertical_priority": "middle",
                    "text_content": ["Main result.", "Comparison summary."],
                    "visual_assets": [],
                },
            ]
        }
    }

    result = TemplateBlockPlanner()(state)

    blocks = result["template_block_plan"]["blocks"]
    assert len(blocks) == 3
    assert len({block["slot_id"] for block in blocks}) == 3
    rewritten_sections = result["story_board"]["spatial_content_plan"]["sections"]
    assert len(rewritten_sections) == 3
    assert all(section["column_assignment"].startswith("slot_") for section in rewritten_sections)
    assert all(section["slot_id"] == section["column_assignment"] for section in rewritten_sections)
    assert all(section.get("capacity_budget") for section in rewritten_sections)


def test_block_capacity_contract_scales_with_slot_area():
    planner = TemplateBlockPlanner()
    state = create_state("/tmp/paper.pdf", width=36, height=51, layout_template="cluster_3_portrait")
    layout = load_block_template_layout("cluster_3_portrait", 36, 51, margin=1.0)
    regions = sorted(layout["regions"], key=lambda item: float(item["area_ratio"]))
    small_region = regions[0]
    large_region = regions[-1]
    sections = [
        {
            "section_id": "large_method",
            "section_title": "Large Method",
            "content_role": "method",
            "region_id": large_region["region_id"],
            "slot_id": large_region["region_id"],
            "visual_assets": [],
        },
        {
            "section_id": "small_note",
            "section_title": "Small Note",
            "content_role": "takeaway",
            "region_id": small_region["region_id"],
            "slot_id": small_region["region_id"],
            "visual_assets": [],
        },
    ]

    contract = planner._build_block_capacity_contract(sections, layout, state)
    by_section = contract["by_section"]

    assert by_section["large_method"]["target_chars"] > by_section["small_note"]["target_chars"]
    assert by_section["large_method"]["target_bullets"] >= by_section["small_note"]["target_bullets"]


def test_block_capacity_contract_reserves_visual_space():
    planner = TemplateBlockPlanner()
    state = create_state("/tmp/paper.pdf", width=36, height=51, layout_template="cluster_3_portrait")
    state["visual_assets"] = {
        "figure_1": {"asset_id": "figure_1", "asset_type": "figure", "aspect": 1.4}
    }
    layout = load_block_template_layout("cluster_3_portrait", 36, 51, margin=1.0)
    region = max(layout["regions"], key=lambda item: float(item["area_ratio"]))
    settings = planner._capacity_settings()
    text_only = {
        "section_id": "text_only",
        "section_title": "Text Only",
        "content_role": "overview",
        "visual_assets": [],
    }
    with_visual = {
        "section_id": "with_visual",
        "section_title": "With Visual",
        "content_role": "method",
        "visual_assets": [{"visual_id": "figure_1"}],
    }

    text_budget = planner._capacity_budget_for_section(text_only, region, state, settings)
    visual_budget = planner._capacity_budget_for_section(with_visual, region, state, settings)

    assert visual_budget["reserved_visual_height"] > text_budget["reserved_visual_height"]
    assert visual_budget["available_text_height"] < text_budget["available_text_height"]
    assert visual_budget["raw_target_chars"] < text_budget["raw_target_chars"]
    assert visual_budget["visual_policy"] in {"reserve_visual_space", "prioritize_visual_scale"}


def test_template_block_planner_capacity_rewrite_preserves_refs_and_expands(tmp_path, monkeypatch):
    class FailingAgent:
        def __init__(self, *args, **kwargs):
            pass

        def step(self, message):
            raise RuntimeError("offline")

    monkeypatch.setattr("src.agents.template_block_planner.LangGraphAgent", FailingAgent)
    state = create_state(str(tmp_path / "paper.pdf"), width=36, height=51, layout_template="cluster_3_portrait")
    state["output_dir"] = str(tmp_path / "output")
    state["resolved_layout_template"] = "cluster_3_portrait"
    state["raw_text"] = (
        "The method uses active geospatial search to choose a sequence of rental units. "
        "The search policy updates after each query and uses property-level information. "
        "The approach accounts for budget and travel-cost constraints. "
        "Experiments compare HAGS with greedy and conventional active search baselines. "
        "Results show improved targeting under uniform and travel-aware budgets."
    )
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "method",
                    "section_title": "Active Search Method",
                    "column_assignment": "middle",
                    "vertical_priority": "top",
                    "text_content": ["Active search selects parcels."],
                    "visual_assets": [],
                    "source_sections": ["method"],
                },
                {
                    "section_id": "results",
                    "section_title": "Results",
                    "column_assignment": "right",
                    "vertical_priority": "middle",
                    "text_content": ["HAGS improves targeting."],
                    "visual_assets": [],
                    "source_sections": ["results"],
                },
                {
                    "section_id": "problem",
                    "section_title": "Problem",
                    "column_assignment": "left",
                    "vertical_priority": "middle",
                    "text_content": ["Outreach teams have limited budgets."],
                    "visual_assets": [],
                    "source_sections": ["problem"],
                },
            ]
        }
    }

    result = TemplateBlockPlanner()(state)
    sections = result["story_board"]["spatial_content_plan"]["sections"]
    method = next(section for section in sections if section["section_id"] == "method")

    assert result["block_capacity_contract"]["blocks"]
    assert result["capacity_planning_report"]["blocks"]
    assert method["slot_id"] == method["column_assignment"]
    assert method["source_sections"] == ["method"]
    assert "capacity_budget" in method
    assert sum(len(item) for item in method["text_content"]) >= len("Active search selects parcels.")


def test_template_block_planner_keypoint_mode_uses_available_unique_slots(tmp_path, monkeypatch):
    class FailingAgent:
        def __init__(self, *args, **kwargs):
            pass

        def step(self, message):
            raise RuntimeError("offline")

    monkeypatch.setattr("src.agents.template_block_planner.LangGraphAgent", FailingAgent)
    template_id = "cluster_62_landscape"
    info = get_block_template_info(template_id)
    canvas = info["recommended_canvas_size"]
    state = create_state(
        str(tmp_path / "paper.pdf"),
        width=canvas["width"],
        height=canvas["height"],
        layout_template=template_id,
    )
    state["output_dir"] = str(tmp_path / "output")
    state["resolved_layout_template"] = template_id
    state["paper_poster_keypoints"] = [
        {"id": index, "key_point": f"Keypoint factual claim {index}", "section": "Method" if index < 7 else "Experiments"}
        for index in range(1, 11)
    ]
    state["poster_reading_order"] = list(range(1, 11))
    state["raw_text"] = " ".join(
        f"Keypoint factual claim {index} is supported by the paper with additional implementation and evaluation context."
        for index in range(1, 11)
    )
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": f"keypoint_{index}",
                    "section_title": f"Point {index}",
                    "column_assignment": "middle",
                    "vertical_priority": "middle",
                    "text_content": [f"Keypoint factual claim {index}."],
                    "visual_assets": [],
                    "keypoint_id": index,
                    "source_section": "Method" if index < 7 else "Experiments",
                    "source_sections": ["Method" if index < 7 else "Experiments"],
                }
                for index in range(1, 11)
            ]
        }
    }

    result = TemplateBlockPlanner()(state)

    blocks = result["template_block_plan"]["blocks"]
    available_slots = len(result["layout_template_metadata"]["regions"])
    assert len(blocks) == available_slots
    assert len({block["slot_id"] for block in blocks}) == available_slots
    assert [block["keypoint_id"] for block in blocks] == list(range(1, available_slots + 1))
    rewritten_sections = result["story_board"]["spatial_content_plan"]["sections"]
    assert [section["keypoint_id"] for section in rewritten_sections] == list(range(1, available_slots + 1))
    assert all(section.get("capacity_budget") for section in rewritten_sections)


def test_template_block_planner_grouped_keypoints_preserve_visuals(tmp_path, monkeypatch):
    class FailingAgent:
        def __init__(self, *args, **kwargs):
            pass

        def step(self, message):
            raise RuntimeError("offline")

    monkeypatch.setattr("src.agents.template_block_planner.LangGraphAgent", FailingAgent)
    template_id = "cluster_104_landscape"
    info = get_block_template_info(template_id)
    canvas = info["recommended_canvas_size"]
    state = create_state(
        str(tmp_path / "paper.pdf"),
        width=canvas["width"],
        height=canvas["height"],
        layout_template=template_id,
    )
    state["output_dir"] = str(tmp_path / "output")
    state["resolved_layout_template"] = template_id
    state["paper_poster_keypoints"] = [
        {"id": index, "key_point": f"Poster keypoint {index}", "section": "Method" if index < 6 else "Results"}
        for index in range(1, 11)
    ]
    state["poster_reading_order"] = list(range(1, 11))
    state["raw_text"] = (
        "am-ELO reformulates Elo scoring with maximum likelihood estimation. "
        "The method models annotator ability and estimates model scores jointly. "
        "Results show lower loss and better prediction performance than baselines. "
        "Robustness tests evaluate perturbations to arena outcomes."
    )
    state["visual_assets"] = {
        "figure_2": {"asset_type": "figure", "aspect": 1.7, "source_path": "/tmp/figure-2.png"},
        "figure_3": {"asset_type": "figure", "aspect": 2.3, "source_path": "/tmp/figure-3.png"},
        "table_3": {"asset_type": "table", "aspect": 0.94, "source_path": "/tmp/table-3.png"},
    }
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "motivation",
                    "section_title": "Motivation",
                    "column_assignment": "left",
                    "vertical_priority": "top",
                    "text_content": ["Arena comparisons can create unstable Elo rankings."],
                    "visual_assets": [],
                    "keypoint_id": 1,
                    "source_keypoint_ids": [1, 2],
                    "source_sections": ["Introduction"],
                    "content_type": "foundation",
                },
                {
                    "section_id": "method",
                    "section_title": "am-ELO Method",
                    "column_assignment": "middle",
                    "vertical_priority": "top",
                    "text_content": ["am-ELO models annotator ability in pairwise comparisons."],
                    "visual_assets": [{"visual_id": "figure_2"}],
                    "keypoint_id": 3,
                    "source_keypoint_ids": [3, 4],
                    "source_sections": ["Method"],
                    "content_type": "method",
                },
                {
                    "section_id": "estimation",
                    "section_title": "MLE Estimate",
                    "column_assignment": "middle",
                    "vertical_priority": "middle",
                    "text_content": ["Likelihood-based estimation avoids sequential update sensitivity."],
                    "visual_assets": [],
                    "keypoint_id": 5,
                    "source_keypoint_ids": [5],
                    "source_sections": ["Method"],
                    "content_type": "method",
                },
                {
                    "section_id": "setup",
                    "section_title": "Experiment Setup",
                    "column_assignment": "left",
                    "vertical_priority": "middle",
                    "text_content": ["Experiments compare Elo, m-ELO, and am-ELO on arena records."],
                    "visual_assets": [],
                    "keypoint_id": 6,
                    "source_keypoint_ids": [6],
                    "source_sections": ["Experiments"],
                    "content_type": "foundation",
                },
                {
                    "section_id": "results",
                    "section_title": "Key Results",
                    "column_assignment": "right",
                    "vertical_priority": "top",
                    "text_content": ["am-ELO obtains lower loss than baseline estimators."],
                    "visual_assets": [{"visual_id": "figure_3"}],
                    "keypoint_id": 7,
                    "source_keypoint_ids": [7],
                    "source_sections": ["Results"],
                    "content_type": "results",
                },
                {
                    "section_id": "table",
                    "section_title": "Prediction",
                    "column_assignment": "right",
                    "vertical_priority": "middle",
                    "text_content": ["Prediction experiments indicate better generalization."],
                    "visual_assets": [{"visual_id": "table_3"}],
                    "keypoint_id": 8,
                    "source_keypoint_ids": [8],
                    "source_sections": ["Results"],
                    "content_type": "results",
                },
                {
                    "section_id": "takeaway",
                    "section_title": "Takeaway",
                    "column_assignment": "right",
                    "vertical_priority": "bottom",
                    "text_content": ["Annotator-aware MLE improves stability for arena-based LLM evaluation."],
                    "visual_assets": [],
                    "keypoint_id": 9,
                    "source_keypoint_ids": [9, 10],
                    "source_sections": ["Robustness", "Conclusion"],
                    "content_type": "takeaway",
                },
            ]
        }
    }

    result = TemplateBlockPlanner()(state)

    sections = result["story_board"]["spatial_content_plan"]["sections"]
    visual_ids = [
        visual["visual_id"]
        for section in sections
        for visual in section.get("visual_assets", [])
    ]
    region_by_id = {
        region["region_id"]: region
        for region in result["layout_template_metadata"]["regions"]
    }
    visual_sections = [section for section in sections if section.get("visual_assets")]
    assert len(sections) == 7
    assert len({section["slot_id"] for section in sections}) == 7
    assert {"figure_2", "figure_3", "table_3"}.issubset(set(visual_ids))
    assert all(section.get("source_keypoint_ids") for section in sections)
    assert all(region_by_id[section["slot_id"]]["can_host_visual"] for section in visual_sections)
    assert all(region_by_id[section["slot_id"]]["text_density_limit"] != "low" for section in visual_sections)


def test_layout_templates_support_block_template_ids():
    template_names = LayoutTemplates.available_template_names()

    assert "cluster_2_landscape" in template_names
    assert "cluster_104_landscape" in template_names
    assert "cluster_3_portrait" in template_names
    info = get_block_template_info("cluster_3_portrait")
    assert info["orientation"] == "portrait"
    layout = LayoutTemplates(36, 51, margin=1.0, col_gap=1.0).get_template("cluster_3_portrait")
    assert layout["layout_mode"] == "template_prior"
    assert layout["orientation"] == "portrait"
    assert [lane["id"] for lane in layout["lanes"]][:2] == ["slot_1", "slot_2"]
    assert "base_layout_template" not in layout


def test_dense_landscape_template_exposes_seven_visual_capable_content_blocks():
    layout = load_block_template_layout("cluster_104_landscape", 54, 27, margin=1.0)
    regions = layout["regions"]

    assert layout["slot_count"] == 7
    assert len(regions) == 7
    assert sum(1 for region in regions if region.get("can_host_visual")) >= 5
    assert any(region["region_id"] == "slot_7" and region.get("text_density_limit") == "low" for region in regions)


def test_template_block_planner_keeps_main_results_out_of_small_low_density_slot():
    layout = load_block_template_layout("cluster_104_landscape", 54, 27, margin=1.0)
    planner = TemplateBlockPlanner()
    region = planner._region_for_ordered_section(
        {
            "section_id": "sec_main_results",
            "section_title": "Main Results",
            "content_role": "results",
            "visual_assets": [],
        },
        layout["regions"],
        {"slot_1", "slot_2", "slot_3", "slot_4", "slot_5"},
        later_visual_count=0,
    )

    assert region["region_id"] == "slot_6"
    slot_7 = next(item for item in layout["regions"] if item["region_id"] == "slot_7")
    assert region["area_ratio"] > slot_7["area_ratio"]
    assert region["can_host_visual"] is True


def test_resolve_poster_dimensions_allows_imported_wide_block_templates():
    width, height = resolve_poster_dimensions("cluster_104_landscape", None, None)

    assert width / height <= 2.1


def test_font_agent_keyword_prompt_uses_narrative_content(monkeypatch):
    captured = {}

    class FakeResponse:
        content = '{"section_keywords": {}, "formatting_summary": {}}'
        input_tokens = 1
        output_tokens = 1

    class FakeAgent:
        def __init__(self, *args, **kwargs):
            pass

        def step(self, message):
            captured["message"] = message
            return FakeResponse()

    state = create_state("/tmp/paper.pdf")
    state["narrative_content"] = {"and": "poster narrative signal"}
    agent = FontAgent()
    monkeypatch.setattr("src.agents.font_agent.LangGraphAgent", FakeAgent)

    agent._identify_keywords({"spatial_content_plan": {"sections": []}}, state)

    assert "poster narrative signal" in captured["message"]


def test_text_cleanup_repairs_mojibake_bullets_and_common_ocr_typos():
    text = "â¢ **Realistic Setting: When costs matter.\\nâ¦ Effcient search improves 42%â70% with î»L_BCE."

    cleaned = normalize_text_for_poster(text)

    assert "â" not in cleaned
    assert "• **Realistic Setting:** When costs matter." in cleaned
    assert "◦ Efficient search improves 42%-70% with lambda L_BCE." in cleaned


def test_text_cleanup_removes_ocr_paths_tables_and_metadata_noise():
    noisy = (
        "A detailed formal presentation follows. ![](_page_4_Figure_0.jpeg) "
        "Figure 2: HAGS policy network architecture. | Search Budget | 15 | 20 |"
    )

    cleaned = normalize_text_for_poster(noisy)

    assert cleaned == "A detailed formal presentation follows."
    assert "_page_" not in cleaned
    assert ".jpeg" not in cleaned
    assert "|" not in cleaned
    assert "Figure 2" not in cleaned


def test_text_cleanup_removes_table_references_and_section_metadata():
    metadata = (
        "HAGS uses a hierarchical policy importance high contains_figures contains_tables "
        "section_name Related Work section_type foundation content."
    )

    assert normalize_text_for_poster(metadata) == "HAGS uses a hierarchical policy"
    assert normalize_text_for_poster("The results are presented in Tables 3 and 4.") == ""
    assert normalize_text_for_poster("Table 5: Main comparison across baselines.") == ""
    assert normalize_text_for_poster("Challenge:** Outreach teams cannot visit every parcel.") == "**Challenge:** Outreach teams cannot visit every parcel."
    assert normalize_text_for_poster("A detailed presentation of the complete method is provided in Algorithm 2 in Supplement.") == ""
    assert normalize_text_for_poster("HAGS is best across budgets, outperforming.") == "HAGS is best across budgets."


def test_renderer_resolves_figure_and_table_paths():
    renderer = Renderer()
    state = create_state("/tmp/paper.pdf")
    state["resolved_visual_assets"] = {
        "methods_figure_1": {"resolved_path": "/tmp/figure-slot.png"},
        "results_table_1": {"resolved_path": "/tmp/table-slot.png"},
    }

    assert renderer._get_resolved_visual_entry("methods_figure_1", "figure_1", state)["resolved_path"] == "/tmp/figure-slot.png"
    assert renderer._get_resolved_visual_entry("results_table_1", "table_1", state)["resolved_path"] == "/tmp/table-slot.png"


def test_visual_asset_agent_disabled_is_slot_preserving_crop_only(tmp_path):
    source_path = tmp_path / "source.png"
    from PIL import Image

    Image.new("RGB", (100, 80), color=(255, 0, 0)).save(source_path)

    state = create_state(str(tmp_path / "paper.pdf"))
    state["output_dir"] = str(tmp_path / "output")
    state["enable_visual_refinement"] = False
    state["visual_assets"] = {
        "figure_1": {
            "asset_id": "figure_1",
            "asset_type": "figure",
            "source_path": str(source_path),
            "resolved_path": None,
            "caption": "Figure 1",
            "aspect": 1.25,
            "provenance": "paper_extracted",
        }
    }
    state["styled_layout"] = [
        {
            "type": "visual",
            "slot_id": "method_figure_1",
            "id": "method_figure_1",
            "visual_id": "figure_1",
            "width": 2.0,
            "height": 1.0,
        }
    ]

    result = VisualAssetAgent()(state)

    assert result["visual_plan"][0]["action"] == "crop_only"
    assert "method_figure_1" in result["resolved_visual_assets"]
    assert Path(result["resolved_visual_assets"]["method_figure_1"]["resolved_path"]).exists()


def test_visual_asset_agent_enabled_generates_for_missing_source_slot(tmp_path):
    state = create_state(str(tmp_path / "paper.pdf"), enable_visual_refinement=True)
    state["output_dir"] = str(tmp_path / "output")
    state["visual_assets"] = {}
    state["styled_layout"] = [
        {
            "type": "visual",
            "slot_id": "method_generated_visual",
            "id": "method_generated_visual",
            "width": 2.0,
            "height": 1.0,
        }
    ]

    result = VisualAssetAgent()(state)

    assert result["visual_plan"][0]["action"] == "generate_new"
    resolved = result["resolved_visual_assets"]["method_generated_visual"]
    assert resolved["asset_type"] == "generated"
    assert Path(resolved["resolved_path"]).exists()


def test_vlm_layout_reviewer_disabled_is_noop():
    state = create_state("/tmp/paper.pdf")
    state["styled_layout"] = [{"type": "title", "id": "title", "x": 1, "y": 1, "width": 10, "height": 2}]

    result = VLMLayoutReviewer()(state)

    assert result.get("vlm_layout_review") is None
    assert result["styled_layout"] == state["styled_layout"]


def test_vlm_layout_reviewer_uses_responses_endpoint(monkeypatch):
    captured = {}

    class FakeResponse:
        headers = {}

        def raise_for_status(self):
            pass

        def json(self):
            return {
                "output_text": '{"overall_score": 90, "accept": true, "issues": [], "patch": [], "visual_asset_recommendations": []}'
            }

    def fake_post(url, headers, json, timeout, stream=False):
        captured["url"] = url
        captured["payload"] = json
        captured["stream"] = stream
        return FakeResponse()

    monkeypatch.setattr("src.agents.vlm_layout_reviewer.requests.post", fake_post)
    reviewer = VLMLayoutReviewer()
    response = reviewer._post_vlm_request(
        "https://example.com/api/v1/responses",
        {"Authorization": "Bearer test"},
        "gpt-5.4",
        "review",
        "data:image/png;base64,abc",
    )
    text = reviewer._extract_response_text(response)

    assert captured["url"] == "https://example.com/api/v1/responses"
    assert captured["payload"]["store"] is False
    assert captured["payload"]["stream"] is True
    assert captured["stream"] is True
    assert captured["payload"]["input"][0]["content"][0]["type"] == "input_text"
    assert captured["payload"]["input"][0]["content"][1]["type"] == "input_image"
    assert "overall_score" in text


def test_vlm_layout_reviewer_falls_back_on_request_failure(tmp_path, monkeypatch):
    state = create_state(str(tmp_path / "paper.pdf"), enable_vlm_layout_review=True)
    state["output_dir"] = str(tmp_path / "output")
    state["poster_preview_path"] = str(tmp_path / "preview.png")
    state["styled_layout"] = [{"type": "title", "id": "title", "x": 1, "y": 1, "width": 10, "height": 2}]
    Path(state["output_dir"], "content").mkdir(parents=True)
    Image.new("RGB", (200, 120), color=(255, 255, 255)).save(state["poster_preview_path"])

    monkeypatch.setenv("VLM_BASE_URL", "https://example.com/v1")
    monkeypatch.setenv("VLM_API_KEY", "test")
    monkeypatch.setenv("VLM_MODEL", "gpt-5.4")

    def fail_post(self, base_url, headers, model, prompt, image_data):
        raise ConnectionResetError("connection reset")

    monkeypatch.setattr(VLMLayoutReviewer, "_post_vlm_request", fail_post)

    result = VLMLayoutReviewer()(state)

    assert result["errors"] == []
    assert result["vlm_layout_review"]["source"] == "fallback"
    assert result["vlm_layout_review"]["accept"] is True
    assert "VLM layout request failed" in result["vlm_layout_review"]["warnings"][0]


def test_gpt_54_uses_openai_chat_provider():
    config = _get_model_config("gpt-5.4")

    assert config.provider == "openai"
    assert config.model_name == "gpt-5.4"


def test_vlm_layout_reviewer_applies_single_safe_patch(tmp_path, monkeypatch):
    state = create_state(str(tmp_path / "paper.pdf"), enable_vlm_layout_review=True)
    state["output_dir"] = str(tmp_path / "output")
    Path(state["output_dir"], "content").mkdir(parents=True)
    preview = tmp_path / "preview.png"
    from PIL import Image

    Image.new("RGB", (200, 120), color=(255, 255, 255)).save(preview)
    state["poster_preview_path"] = str(preview)
    state["layout_template_metadata"] = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0).get_template(
        "three_column_postergen",
        header_height=(36 - 2) * 0.18,
    )
    left_lane = state["layout_template_metadata"]["lanes"][0]
    state["styled_layout"] = [
        {
            "type": "section_container",
            "section_id": "intro",
            "lane_id": "left",
            "x": left_lane["x"],
            "y": left_lane["y"],
            "width": left_lane["w"],
            "height": 4.0,
            "priority": 0.1,
        },
        {
            "type": "visual",
            "id": "intro_visual",
            "slot_id": "intro_visual",
            "visual_id": "figure_1",
            "x": left_lane["x"] + 1.0,
            "y": left_lane["y"] + 1.0,
            "width": 4.0,
            "height": 2.0,
            "priority": 0.4,
        },
    ]

    def fake_review(self, state):
        return {
            "overall_score": 82,
            "accept": False,
            "issues": [{"severity": "medium", "category": "whitespace", "target": "intro_visual"}],
            "patch": [{"target": "intro_visual", "op": "increase_visual_scale", "value": 1.1}],
            "visual_asset_recommendations": [],
        }

    monkeypatch.setattr(VLMLayoutReviewer, "_review_or_fallback", fake_review)
    result = VLMLayoutReviewer()(state)

    visual = next(element for element in result["styled_layout"] if element.get("id") == "intro_visual")
    assert result["vlm_patch_applied"] is True
    assert result["vlm_reflow_required"] is True
    assert result["vlm_review_count"] == 1
    assert visual["width"] > 4.0


def test_vlm_layout_reviewer_fast_mode_records_patch_without_reflow(tmp_path, monkeypatch):
    state = create_state(str(tmp_path / "paper.pdf"), enable_vlm_layout_review=True)
    state["output_dir"] = str(tmp_path / "output")
    state["template_fast_mode"] = True
    state["template_layout_mode"] = "template_prior"
    state["styled_layout"] = [{"type": "title", "id": "title", "x": 1, "y": 1, "width": 10, "height": 2}]
    Path(state["output_dir"], "content").mkdir(parents=True)

    reviewer = VLMLayoutReviewer()
    monkeypatch.setattr(
        reviewer,
        "_review_or_fallback",
        lambda _state: {
            "source": "test",
            "overall_score": 62,
            "accept": False,
            "issues": [{"severity": "medium", "category": "whitespace", "description": "underfilled"}],
            "patch": [{"operation": "move", "target": "title", "dx": 0.1, "dy": 0.0}],
            "warnings": [],
        },
    )

    result = reviewer(state)

    assert result["vlm_reflow_required"] is False
    assert result["vlm_patch_applied"] is False
    assert result["template_repair_required"] is False
    assert result["vlm_layout_patch"]
    assert "did not apply" in result["vlm_layout_review"]["warnings"][0]


def test_visual_legibility_fast_mode_is_report_only(tmp_path, monkeypatch):
    state = create_state(str(tmp_path / "paper.pdf"), enable_visual_legibility_review=True)
    state["output_dir"] = str(tmp_path / "output")
    state["template_fast_mode"] = True
    state["template_layout_mode"] = "template_prior"
    Path(state["output_dir"], "content").mkdir(parents=True)

    reviewer = VisualLegibilityReviewer()
    monkeypatch.setattr(
        reviewer,
        "_review_or_fallback",
        lambda _state: {
            "needs_relayout": True,
            "issues": [{"severity": "medium", "target": "slot_3", "description": "small figure text"}],
            "layout_recommendation": {"target_region": "slot_3", "action": "promote_region", "reason": "small"},
            "warnings": [],
        },
    )
    monkeypatch.setattr(reviewer, "_merge_heuristic_review", lambda _state, review: review)

    result = reviewer(state)

    assert result["template_repair_required"] is False
    assert result["adaptive_relayout_required"] is False
    assert "does not trigger automatic relayout" in result["visual_legibility_review"]["warnings"][0]


def test_layout_templates_expose_multiple_geometry_families():
    templates = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0)

    three_col = templates.get_template("three_column_postergen", header_height=6.0)
    two_plus_one = templates.get_template("two_plus_one_mixed", header_height=6.0)
    one_plus_two = templates.get_template("one_plus_two_mixed", header_height=6.0)
    single_col = templates.get_template("single_column_vertical", header_height=6.0)
    adaptive = templates.get_template(
        "adaptive_three_column",
        header_height=6.0,
        width_ratios={"left": 0.85, "middle": 1.30, "right": 0.85},
    )

    assert len(three_col["lanes"]) == 3
    assert round(three_col["lanes"][0]["w"], 4) == round(three_col["lanes"][1]["w"], 4)

    assert two_plus_one["lanes"][2]["w"] > two_plus_one["lanes"][0]["w"]
    assert one_plus_two["lanes"][0]["w"] > one_plus_two["lanes"][1]["w"]

    assert len({round(lane["x"], 4) for lane in single_col["lanes"]}) == 1
    assert single_col["lanes"][0]["y"] < single_col["lanes"][1]["y"] < single_col["lanes"][2]["y"]

    assert adaptive["template_name"] == "adaptive_three_column"
    assert adaptive["lanes"][1]["w"] > adaptive["lanes"][0]["w"]
    assert adaptive["lanes"][1]["w"] > adaptive["lanes"][2]["w"]


def test_template_extractor_builds_expected_template_schema():
    image_path = Path("template/poster(1).png")
    if not image_path.exists():
        return

    template, raw = build_template(image_path)

    assert template["geometry_policy"] == "soft"
    assert template["source_lanes"]
    assert template["panel_style_tokens"]
    assert len(template["lanes"]) == 3
    assert raw["ocr_blocks"]
    for box in [template["header"], *template["lanes"], *template["logo_regions"]]:
        assert 0 <= box["x"] <= 1
        assert 0 <= box["y"] <= 1
        assert 0 < box["w"] <= 1
        assert 0 < box["h"] <= 1


def test_extracted_template_registry_and_scaling():
    template_ids = set(list_extracted_template_ids())
    if not template_ids:
        return

    assert "extracted_poster1_landscape_three_panel" in template_ids
    assert "extracted_poster2_landscape_multi_panel" in template_ids
    assert "extracted_poster3_portrait_section_band" in template_ids

    loaded = load_extracted_template("extracted_poster3_portrait_section_band")
    assert loaded is not None
    assert loaded["orientation"] == "portrait"

    layout = LayoutTemplates(36, 54, margin=1.0, col_gap=1.0).get_template(
        "extracted_poster3_portrait_section_band"
    )
    assert layout["template_name"] == "extracted_poster3_portrait_section_band"
    assert layout["orientation"] == "portrait"
    assert layout["geometry_policy"] == "soft"
    assert len(layout["lanes"]) == 3
    assert layout["lanes"][0]["w"] > layout["lanes"][0]["h"]
    assert layout["lanes"][0]["y"] < layout["lanes"][1]["y"] < layout["lanes"][2]["y"]
    assert layout["source_lanes"][0]["h"] != layout["lanes"][0]["h"]

    horizontal = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0).get_template(
        "extracted_poster1_landscape_three_panel",
        header_height=6.0,
    )
    body_height = 36 - 1.0 - horizontal["lanes"][0]["y"]
    assert len(horizontal["lanes"]) == 3
    assert horizontal["lanes"][0]["h"] >= body_height * 0.95
    assert horizontal["lanes"][0]["y"] < horizontal["source_lanes"][0]["y"]


def test_layout_agent_resolves_extracted_template_metadata():
    if "extracted_poster3_portrait_section_band" not in set(list_extracted_template_ids()):
        return

    state = create_state(
        "/tmp/paper.pdf",
        width=36,
        height=54,
        layout_template="extracted_poster3_portrait_section_band",
    )

    template = LayoutAgent()._resolve_template_layout(state)

    assert template["template_name"] == "extracted_poster3_portrait_section_band"
    assert state["resolved_layout_template"] == "extracted_poster3_portrait_section_band"
    assert template["header"]["w"] == 36 * 0.95
    assert template["lanes"][0]["w"] == 34


def test_layout_agent_respects_requested_template_geometry():
    state = create_state("/tmp/paper.pdf", layout_template="one_plus_two_mixed")
    state["narrative_content"] = {"meta": {"poster_title": "Paper", "authors": "Authors"}}
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "intro",
                    "section_title": "Intro",
                    "column_assignment": "left",
                    "vertical_priority": "top",
                    "text_content": ["Point A", "Point B"],
                    "visual_assets": [],
                    "importance_level": 2,
                },
                {
                    "section_id": "method",
                    "section_title": "Method",
                    "column_assignment": "middle",
                    "vertical_priority": "top",
                    "text_content": ["Point A", "Point B"],
                    "visual_assets": [],
                    "importance_level": 1,
                },
                {
                    "section_id": "results",
                    "section_title": "Results",
                    "column_assignment": "right",
                    "vertical_priority": "top",
                    "text_content": ["Point A", "Point B"],
                    "visual_assets": [],
                    "importance_level": 2,
                },
            ]
        }
    }

    result = LayoutAgent()(state, mode="initial")
    sections = {
        element["section_id"]: element
        for element in result["initial_layout_data"]
        if element.get("type") == "section_container"
    }

    assert sections["intro"]["width"] > sections["method"]["width"]
    assert sections["intro"]["width"] > sections["results"]["width"]


def test_layout_agent_respects_adaptive_lane_widths():
    state = create_state("/tmp/paper.pdf", layout_template="three_column_postergen")
    state["adaptive_lane_widths"] = {"left": 0.85, "middle": 1.30, "right": 0.85}
    state["narrative_content"] = {"meta": {"poster_title": "Paper", "authors": "Authors"}}
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "intro",
                    "section_title": "Intro",
                    "column_assignment": "left",
                    "vertical_priority": "top",
                    "text_content": ["Point A"],
                    "visual_assets": [],
                    "importance_level": 2,
                },
                {
                    "section_id": "method",
                    "section_title": "Method",
                    "column_assignment": "middle",
                    "vertical_priority": "top",
                    "text_content": ["Point A"],
                    "visual_assets": [],
                    "importance_level": 1,
                },
                {
                    "section_id": "results",
                    "section_title": "Results",
                    "column_assignment": "right",
                    "vertical_priority": "top",
                    "text_content": ["Point A"],
                    "visual_assets": [],
                    "importance_level": 2,
                },
            ]
        }
    }

    result = LayoutAgent()(state, mode="initial")
    sections = {
        element["section_id"]: element
        for element in result["initial_layout_data"]
        if element.get("type") == "section_container"
    }

    assert result["resolved_layout_template"] == "adaptive_three_column"
    assert sections["method"]["width"] > sections["intro"]["width"]
    assert sections["method"]["width"] > sections["results"]["width"]


def test_visual_legibility_heuristic_requests_middle_lane_for_wide_visual():
    state = create_state("/tmp/paper.pdf", enable_visual_legibility_review=True, enable_adaptive_column_width=True)
    state["layout_template_metadata"] = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0).get_template(
        "three_column_postergen",
        header_height=6.0,
    )
    middle_lane = state["layout_template_metadata"]["lanes"][1]
    state["visual_assets"] = {
        "figure_2": {
            "asset_id": "figure_2",
            "caption": "Hierarchical method pipeline overview",
            "asset_type": "figure",
        }
    }
    state["styled_layout"] = [
        {
            "type": "visual",
            "id": "method_figure",
            "slot_id": "method_figure",
            "visual_id": "figure_2",
            "x": middle_lane["x"] + 0.3,
            "y": middle_lane["y"] + 1.0,
            "width": 15.0,
            "height": 5.0,
        }
    ]

    review = VisualLegibilityReviewer()._heuristic_review(state)

    assert review["needs_relayout"] is True
    assert review["layout_recommendation"]["target_lane"] == "middle"


def test_adaptive_column_relayout_sets_template_and_saves_decision(tmp_path):
    state = create_state(
        str(tmp_path / "paper.pdf"),
        enable_visual_legibility_review=True,
        enable_adaptive_column_width=True,
    )
    state["output_dir"] = str(tmp_path / "output")
    state["adaptive_relayout_required"] = True
    state["visual_legibility_review"] = {
        "needs_relayout": True,
        "layout_recommendation": {
            "target_lane": "middle",
            "action": "widen_lane",
            "preferred_width_ratio": 1.3,
            "reason": "Middle visual text is too small.",
        },
        "issues": [],
    }

    result = AdaptiveColumnRelayoutAgent()(state)

    assert result["layout_template"] == "adaptive_three_column"
    assert result["adaptive_relayout_count"] == 1
    assert result["adaptive_lane_widths"]["middle"] > result["adaptive_lane_widths"]["left"]
    assert Path(state["output_dir"], "content", "adaptive_layout_decision.json").exists()


def test_layout_agent_single_column_template_stacks_semantic_lanes():
    state = create_state("/tmp/paper.pdf", layout_template="single_column_vertical")
    state["narrative_content"] = {"meta": {"poster_title": "Paper", "authors": "Authors"}}
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "intro",
                    "section_title": "Intro",
                    "column_assignment": "left",
                    "vertical_priority": "top",
                    "text_content": ["Point A", "Point B"],
                    "visual_assets": [],
                    "importance_level": 2,
                },
                {
                    "section_id": "method",
                    "section_title": "Method",
                    "column_assignment": "middle",
                    "vertical_priority": "top",
                    "text_content": ["Point A", "Point B"],
                    "visual_assets": [],
                    "importance_level": 1,
                },
                {
                    "section_id": "results",
                    "section_title": "Results",
                    "column_assignment": "right",
                    "vertical_priority": "top",
                    "text_content": ["Point A", "Point B"],
                    "visual_assets": [],
                    "importance_level": 2,
                },
            ]
        }
    }

    result = LayoutAgent()(state, mode="initial")
    sections = {
        element["section_id"]: element
        for element in result["initial_layout_data"]
        if element.get("type") == "section_container"
    }

    assert sections["intro"]["x"] == sections["method"]["x"] == sections["results"]["x"]
    assert sections["intro"]["y"] < sections["method"]["y"] < sections["results"]["y"]


def test_micro_layout_refiner_packs_overflowing_lane_without_lane_overflow():
    state = create_state("/tmp/paper.pdf", layout_template="three_column_postergen")
    state["layout_template_metadata"] = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0).get_template(
        "three_column_postergen",
        header_height=(36 - 2) * 0.18,
    )
    state["styled_layout"] = [
        {
            "type": "section_container",
            "section_id": "s1",
            "lane_id": "left",
            "x": 1.0,
            "y": 7.12,
            "width": 16.66,
            "height": 12.0,
            "importance_level": 2,
            "priority": 0.1,
        },
        {
            "type": "text",
            "id": "s1_text",
            "x": 1.3,
            "y": 8.0,
            "width": 16.06,
            "height": 11.0,
            "content": "Point A\nPoint B\nPoint C\nPoint D\nPoint E\nPoint F",
            "font_family": "Arial",
            "font_size": 44,
            "font_color": "#000000",
            "priority": 0.5,
        },
        {
            "type": "section_container",
            "section_id": "s2",
            "lane_id": "left",
            "x": 1.0,
            "y": 20.5,
            "width": 16.66,
            "height": 12.0,
            "importance_level": 2,
            "priority": 0.1,
        },
        {
            "type": "text",
            "id": "s2_text",
            "x": 1.3,
            "y": 21.3,
            "width": 16.06,
            "height": 10.8,
            "content": "Point A\nPoint B\nPoint C\nPoint D\nPoint E\nPoint F",
            "font_family": "Arial",
            "font_size": 44,
            "font_color": "#000000",
            "priority": 0.5,
        },
    ]

    result = MicroLayoutRefiner()(state)
    left_lane = state["layout_template_metadata"]["lanes"][0]
    left_sections = [
        element for element in result["styled_layout"]
        if element.get("type") == "section_container" and element.get("lane_id") == "left"
    ]

    assert left_sections
    assert max(section["y"] + section["height"] for section in left_sections) <= left_lane["y"] + left_lane["h"] + 0.05


def test_micro_layout_refiner_expands_underfilled_lane():
    state = create_state("/tmp/paper.pdf", layout_template="three_column_postergen")
    state["layout_template_metadata"] = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0).get_template(
        "three_column_postergen",
        header_height=(36 - 2) * 0.18,
    )
    left_lane = state["layout_template_metadata"]["lanes"][0]
    state["styled_layout"] = [
        {
            "type": "section_container",
            "section_id": "s1",
            "lane_id": "left",
            "x": left_lane["x"],
            "y": left_lane["y"],
            "width": left_lane["w"],
            "height": 4.0,
            "importance_level": 2,
            "priority": 0.1,
        },
        {
            "type": "section_title",
            "id": "s1_title",
            "x": left_lane["x"] + 0.2,
            "y": left_lane["y"] + 0.2,
            "width": left_lane["w"] - 0.4,
            "height": 1.0,
            "font_size": 64,
            "priority": 0.2,
        },
        {
            "type": "text",
            "id": "s1_text",
            "x": left_lane["x"] + 0.3,
            "y": left_lane["y"] + 1.4,
            "width": left_lane["w"] - 0.6,
            "height": 2.0,
            "content": "Point A\nPoint B\nPoint C",
            "font_family": "Arial",
            "font_size": 44,
            "font_color": "#000000",
            "priority": 0.5,
        },
        {
            "type": "section_container",
            "section_id": "s2",
            "lane_id": "left",
            "x": left_lane["x"],
            "y": left_lane["y"] + 5.0,
            "width": left_lane["w"],
            "height": 4.0,
            "importance_level": 2,
            "priority": 0.1,
        },
        {
            "type": "section_title",
            "id": "s2_title",
            "x": left_lane["x"] + 0.2,
            "y": left_lane["y"] + 5.2,
            "width": left_lane["w"] - 0.4,
            "height": 1.0,
            "font_size": 64,
            "priority": 0.2,
        },
        {
            "type": "text",
            "id": "s2_text",
            "x": left_lane["x"] + 0.3,
            "y": left_lane["y"] + 6.4,
            "width": left_lane["w"] - 0.6,
            "height": 2.0,
            "content": "Point D\nPoint E\nPoint F",
            "font_family": "Arial",
            "font_size": 44,
            "font_color": "#000000",
            "priority": 0.5,
        },
    ]

    original_bottom = max(
        element["y"] + element["height"]
        for element in state["styled_layout"]
        if element.get("type") == "section_container"
    )
    result = MicroLayoutRefiner()(state)
    refined_sections = [
        element for element in result["styled_layout"]
        if element.get("type") == "section_container" and element.get("lane_id") == "left"
    ]
    refined_bottom = max(section["y"] + section["height"] for section in refined_sections)

    assert refined_bottom > original_bottom
    assert refined_bottom <= left_lane["y"] + left_lane["h"] + 0.05


def test_micro_layout_refiner_validation_rejects_child_outside_container():
    state = create_state("/tmp/paper.pdf", layout_template="three_column_postergen")
    template = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0).get_template(
        "three_column_postergen",
        header_height=(36 - 2) * 0.18,
    )
    lane_map = {lane["id"]: lane for lane in template["lanes"]}
    lane = lane_map["left"]
    layout = [
        {
            "type": "section_container",
            "section_id": "hags_core",
            "lane_id": "left",
            "x": lane["x"],
            "y": lane["y"],
            "width": lane["w"],
            "height": 3.0,
        },
        {
            "type": "text",
            "id": "hags_core_text",
            "x": lane["x"] + 0.3,
            "y": lane["y"] + 1.0,
            "width": lane["w"] - 0.6,
            "height": 3.5,
        },
    ]

    validation = MicroLayoutRefiner()._validate_refined_layout(layout, lane_map, state)

    assert any("child vertical overflow" in issue for issue in validation["issues"])


def test_vlm_layout_reviewer_syncs_container_after_patch():
    state = create_state("/tmp/paper.pdf", layout_template="three_column_postergen", enable_vlm_layout_review=True)
    template = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0).get_template(
        "three_column_postergen",
        header_height=(36 - 2) * 0.18,
    )
    state["layout_template_metadata"] = template
    lane = template["lanes"][0]
    layout = [
        {
            "type": "section_container",
            "section_id": "hags_core",
            "lane_id": "left",
            "x": lane["x"],
            "y": lane["y"],
            "width": lane["w"],
            "height": 4.0,
        },
        {
            "type": "text",
            "id": "hags_core_text",
            "x": lane["x"] + 0.3,
            "y": lane["y"] + 1.0,
            "width": lane["w"] - 0.6,
            "height": 3.8,
            "font_size": 44,
        },
    ]
    patch = [{"target": "hags_core_text", "op": "increase_font_size", "value": 2}]

    patched = VLMLayoutReviewer()._apply_safe_patch(layout, patch, state)
    container = next(element for element in patched if element.get("type") == "section_container")
    text = next(element for element in patched if element.get("id") == "hags_core_text")

    assert patched is not None
    assert container["y"] + container["height"] >= text["y"] + text["height"]


def test_micro_layout_refiner_handles_two_plus_one_mixed_without_right_lane_overflow():
    state = create_state("/tmp/paper.pdf", layout_template="two_plus_one_mixed")
    state["layout_template_metadata"] = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0).get_template(
        "two_plus_one_mixed",
        header_height=(36 - 2) * 0.18,
    )
    right_lane = state["layout_template_metadata"]["lanes"][2]
    state["styled_layout"] = [
        {
            "type": "section_container",
            "section_id": "results",
            "lane_id": "right",
            "x": right_lane["x"],
            "y": right_lane["y"],
            "width": right_lane["w"],
            "height": right_lane["h"] * 0.75,
            "importance_level": 1,
            "priority": 0.1,
        },
        {
            "type": "section_title",
            "id": "results_title",
            "x": right_lane["x"] + 0.2,
            "y": right_lane["y"] + 0.2,
            "width": right_lane["w"] - 0.4,
            "height": 1.0,
            "font_size": 64,
            "priority": 0.2,
        },
        {
            "type": "visual",
            "id": "results_visual",
            "slot_id": "results_visual",
            "visual_id": "table_1",
            "x": right_lane["x"] + 0.4,
            "y": right_lane["y"] + 1.6,
            "width": right_lane["w"] - 0.8,
            "height": 7.5,
            "priority": 0.4,
        },
        {
            "type": "text",
            "id": "results_text",
            "x": right_lane["x"] + 0.3,
            "y": right_lane["y"] + 9.6,
            "width": right_lane["w"] - 0.6,
            "height": 15.0,
            "content": "\n".join([f"Result point {i}" for i in range(1, 20)]),
            "font_family": "Arial",
            "font_size": 44,
            "font_color": "#000000",
            "priority": 0.5,
        },
    ]

    result = MicroLayoutRefiner()(state)
    right_sections = [
        element for element in result["styled_layout"]
        if element.get("type") == "section_container" and element.get("lane_id") == "right"
    ]

    assert right_sections
    assert max(section["y"] + section["height"] for section in right_sections) <= right_lane["y"] + right_lane["h"] + 0.05


def test_micro_layout_refiner_handles_one_plus_two_mixed_without_left_lane_overflow():
    state = create_state("/tmp/paper.pdf", layout_template="one_plus_two_mixed")
    state["layout_template_metadata"] = LayoutTemplates(54, 36, margin=1.0, col_gap=1.0).get_template(
        "one_plus_two_mixed",
        header_height=(36 - 2) * 0.18,
    )
    left_lane = state["layout_template_metadata"]["lanes"][0]
    state["styled_layout"] = [
        {
            "type": "section_container",
            "section_id": "intro",
            "lane_id": "left",
            "x": left_lane["x"],
            "y": left_lane["y"],
            "width": left_lane["w"],
            "height": left_lane["h"] * 0.82,
            "importance_level": 2,
            "priority": 0.1,
        },
        {
            "type": "section_title",
            "id": "intro_title",
            "x": left_lane["x"] + 0.2,
            "y": left_lane["y"] + 0.2,
            "width": left_lane["w"] - 0.4,
            "height": 1.0,
            "font_size": 64,
            "priority": 0.2,
        },
        {
            "type": "text",
            "id": "intro_text",
            "x": left_lane["x"] + 0.3,
            "y": left_lane["y"] + 1.4,
            "width": left_lane["w"] - 0.6,
            "height": 18.0,
            "content": "\n".join([f"Background point {i}" for i in range(1, 26)]),
            "font_family": "Arial",
            "font_size": 44,
            "font_color": "#000000",
            "priority": 0.5,
        },
        {
            "type": "visual",
            "id": "intro_visual",
            "slot_id": "intro_visual",
            "visual_id": "figure_1",
            "x": left_lane["x"] + 0.5,
            "y": left_lane["y"] + 20.0,
            "width": left_lane["w"] - 1.0,
            "height": 6.5,
            "priority": 0.4,
        },
    ]

    result = MicroLayoutRefiner()(state)
    left_sections = [
        element for element in result["styled_layout"]
        if element.get("type") == "section_container" and element.get("lane_id") == "left"
    ]

    assert left_sections
    assert max(section["y"] + section["height"] for section in left_sections) <= left_lane["y"] + left_lane["h"] + 0.05


def test_template_selector_prefers_balanced_three_column_layout():
    selector = TemplateSelector(load_config())
    state = create_state("/tmp/paper.pdf", layout_template="adaptive_auto")
    state["visual_assets"] = {
        "figure_1": {"aspect": 2.0, "asset_type": "figure"},
        "figure_2": {"aspect": 2.1, "asset_type": "figure"},
        "table_1": {"aspect": 3.0, "asset_type": "table"},
        "table_2": {"aspect": 3.2, "asset_type": "table"},
    }
    structured_sections = {
        "paper_sections": [
            {"section_name": "Introduction", "section_type": "foundation", "key_points": ["A", "B"], "contains_figures": ["figure_1"], "contains_tables": []},
            {"section_name": "Method", "section_type": "method", "key_points": ["A", "B"], "contains_figures": ["figure_2"], "contains_tables": []},
            {"section_name": "Results", "section_type": "evaluation", "key_points": ["A", "B"], "contains_figures": [], "contains_tables": ["table_1", "table_2"]},
        ]
    }
    classified_visuals = {
        "key_visual": "figure_2",
        "problem_illustration": ["figure_1"],
        "method_workflow": ["figure_2"],
        "main_results": ["table_1"],
        "comparative_results": ["table_2"],
        "supporting": [],
    }

    result = selector.select(state, structured_sections, classified_visuals, state["visual_assets"])

    assert result["selection_mode"] == "adaptive_auto"
    assert result["selected_template"] in {"three_column_postergen", "two_plus_one_mixed", "one_plus_two_mixed"}
    assert result["selected_template"] == "three_column_postergen"


def test_adaptive_auto_layout_uses_dense_column_gaps():
    cfg = load_config()
    state = create_state("/tmp/paper.pdf", layout_template="adaptive_auto", width=54, height=36)
    state["resolved_layout_template"] = "three_column_postergen"
    layout = LayoutAgent()._resolve_template_layout(state)
    lanes = layout["lanes"]
    gaps = [
        round(float(lanes[index + 1]["x"]) - (float(lanes[index]["x"]) + float(lanes[index]["w"])), 4)
        for index in range(len(lanes) - 1)
    ]

    assert layout["template_name"] == "three_column_postergen"
    assert max(gaps) <= cfg["adaptive_auto_dense_layout"]["max_column_gap_inches"]
    assert lanes[0]["x"] == pytest.approx(cfg["adaptive_auto_dense_layout"]["poster_margin"])


def test_template_selector_exposes_right_heavy_lane_preference_in_auto_mode():
    selector = TemplateSelector(load_config())
    state = create_state("/tmp/paper.pdf", layout_template="adaptive_auto")
    state["visual_assets"] = {
        "figure_1": {"aspect": 2.4, "asset_type": "figure"},
        "figure_2": {"aspect": 2.4, "asset_type": "figure"},
        "table_1": {"aspect": 1.6, "asset_type": "table"},
        "table_2": {"aspect": 1.7, "asset_type": "table"},
    }
    structured_sections = {
        "paper_sections": [
            {"section_name": "Introduction", "section_type": "foundation", "key_points": ["A"], "contains_figures": ["figure_1"], "contains_tables": []},
            {"section_name": "Method", "section_type": "method", "key_points": ["A"], "contains_figures": ["figure_2"], "contains_tables": []},
            {
                "section_name": "Results",
                "section_type": "evaluation",
                "key_points": ["A", "B", "C", "D", "E", "F", "G"],
                "contains_figures": [],
                "contains_tables": ["table_1", "table_2"],
            },
        ]
    }
    classified_visuals = {
        "key_visual": "figure_2",
        "problem_illustration": ["figure_1"],
        "method_workflow": ["figure_2"],
        "main_results": ["table_1"],
        "comparative_results": ["table_2"],
        "supporting": [],
    }

    result = selector.select(state, structured_sections, classified_visuals, state["visual_assets"])

    assert result["preferred_template"] == "two_plus_one_mixed"
    assert result["selected_template"] in {"three_column_postergen", "two_plus_one_mixed", "one_plus_two_mixed"}


def test_template_selector_respects_manual_template_request():
    selector = TemplateSelector(load_config())
    state = create_state("/tmp/paper.pdf", layout_template="one_plus_two_mixed")

    result = selector.select(
        state,
        structured_sections={"paper_sections": []},
        classified_visuals={},
        visual_assets={},
    )

    assert result["selection_mode"] == "manual"
    assert result["selected_template"] == "one_plus_two_mixed"


def test_template_selector_auto_uses_standard_template_whitelist():
    selector = TemplateSelector(load_config())
    state = create_state("/tmp/paper.pdf", layout_template="auto")
    state["paper_poster_keypoints"] = [
        {"id": index, "key_point": f"Keypoint {index}", "section": "Method"}
        for index in range(1, 11)
    ]

    result = selector.select(
        state,
        structured_sections={"paper_sections": []},
        classified_visuals={},
        visual_assets={},
    )

    assert result["selection_mode"] == "standard_auto"
    expected_templates = set(load_config()["standard_template_policy"]["auto_templates"])
    assert result["selected_template"] in expected_templates
    selected = next(candidate for candidate in result["candidates"] if candidate["template_name"] == result["selected_template"])
    assert selected["template_name"] in expected_templates


def test_template_selector_dense_standard_auto_prefers_dense_landscape_template():
    selector = TemplateSelector(load_config())
    state = create_state("/tmp/paper.pdf", layout_template="auto", width=54, height=27)
    visual_assets = {
        "figure_1": {"asset_type": "figure", "aspect": 2.0},
        "figure_2": {"asset_type": "figure", "aspect": 1.5},
        "table_1": {"asset_type": "table", "aspect": 2.4},
    }
    structured_sections = {
        "paper_sections": [
            {"section_name": f"Section {index}", "section_type": "method", "key_points": ["A"]}
            for index in range(6)
        ]
    }

    result = selector.select(
        state,
        structured_sections=structured_sections,
        classified_visuals={"main_results": ["table_1"]},
        visual_assets=visual_assets,
    )

    assert result["selection_mode"] == "standard_auto"
    assert result["selected_template"] == "cluster_104_landscape"


def _block_refinement_state(tmp_path, utilization=0.45):
    state = create_state(
        str(tmp_path / "paper.pdf"),
        layout_template="cluster_3_portrait",
        enable_block_vlm_review=True,
    )
    state["output_dir"] = str(tmp_path / "output")
    lane = {"id": "slot_1", "x": 1.0, "y": 2.0, "w": 8.0, "h": 30.0}
    used_height = lane["h"] * utilization
    state["layout_template_metadata"] = {
        "template_name": "cluster_3_portrait",
        "layout_mode": "template_prior",
        "lanes": [lane],
    }
    state["story_board"] = {
        "spatial_content_plan": {
            "sections": [
                {
                    "section_id": "method",
                    "section_title": "Method",
                    "column_assignment": "slot_1",
                    "slot_id": "slot_1",
                    "vertical_priority": "top",
                    "text_content": ["Existing factual bullet about the method."],
                    "visual_assets": [{"visual_id": "figure_1"}],
                }
            ]
        }
    }
    state["styled_layout"] = [
        {
            "type": "section_container",
            "section_id": "method",
            "lane_id": "slot_1",
            "slot_id": "slot_1",
            "x": lane["x"],
            "y": lane["y"],
            "width": lane["w"],
            "height": used_height,
        },
        {
            "type": "text",
            "id": "method_text",
            "section_id": "method",
            "lane_id": "slot_1",
            "slot_id": "slot_1",
            "x": lane["x"] + 0.3,
            "y": lane["y"] + 1.0,
            "width": lane["w"] - 0.6,
            "height": max(used_height - 1.0, 0.5),
            "font_size": 44,
            "line_spacing": 1.0,
            "content": "Existing factual bullet about the method.",
        },
    ]
    return state


def test_block_occupancy_analyzer_formula_actions(tmp_path):
    analyzer = BlockOccupancyAnalyzer()

    low = analyzer.analyze(_block_refinement_state(tmp_path, utilization=0.45))["blocks"][0]
    moderate = analyzer.analyze(_block_refinement_state(tmp_path, utilization=0.92))["blocks"][0]
    near_target = analyzer.analyze(_block_refinement_state(tmp_path, utilization=0.96))["blocks"][0]
    crowded = analyzer.analyze(_block_refinement_state(tmp_path, utilization=0.99))["blocks"][0]

    assert low["action"] == "expand"
    assert low["target_extra_chars"] > moderate["target_extra_chars"]
    assert moderate["action"] == "expand"
    assert 0 < moderate["target_extra_chars"] < low["target_extra_chars"]
    assert near_target["action"] == "keep"
    assert near_target["target_extra_chars"] == 0
    assert crowded["action"] == "reduce"


def test_block_occupancy_analyzer_uses_real_child_content_not_container_background(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=1.0)
    state["styled_layout"][0]["height"] = 30.0
    state["styled_layout"][1]["y"] = 3.0
    state["styled_layout"][1]["height"] = 8.0

    block = BlockOccupancyAnalyzer().analyze(state)["blocks"][0]

    assert block["container_bbox"]["h"] == pytest.approx(30.0)
    assert block["used_height"] == pytest.approx(9.0)
    assert block["bottom_whitespace"] == pytest.approx(21.0)
    assert block["action"] == "expand"
    assert block["target_extra_chars"] > 0


def test_final_quality_gate_rejects_block_below_hard_minimum(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=0.50)
    state["output_dir"] = str(tmp_path / "output")
    state["template_layout_mode"] = "template_prior"
    state["final_poster_accepted"] = True
    content_dir = Path(state["output_dir"]) / "content"
    content_dir.mkdir(parents=True, exist_ok=True)
    (content_dir / "micro_layout_report.json").write_text(
        json.dumps({"validation": {"issues": []}}),
        encoding="utf-8",
    )

    result = _run_final_quality_gate(state)

    assert result["final_poster_accepted"] is False
    assert result["final_quality_gate"]["accepted"] is False
    assert result["final_quality_gate"]["failures"][0]["category"] == "occupancy"
    assert (content_dir / "final_quality_gate.json").exists()


def test_final_quality_gate_rejects_mean_below_96_percent(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=0.95)
    state["output_dir"] = str(tmp_path / "output")
    state["template_layout_mode"] = "template_prior"
    state["final_poster_accepted"] = True
    content_dir = Path(state["output_dir"]) / "content"
    content_dir.mkdir(parents=True, exist_ok=True)
    (content_dir / "micro_layout_report.json").write_text(
        json.dumps({"validation": {"issues": []}}),
        encoding="utf-8",
    )

    result = _run_final_quality_gate(state)

    assert result["final_poster_accepted"] is False
    assert result["final_quality_gate"]["accepted"] is False
    assert any(
        failure["category"] == "occupancy_mean"
        for failure in result["final_quality_gate"]["failures"]
    )


def test_block_content_refiner_expands_underfilled_block_without_changing_refs(tmp_path, monkeypatch):
    state = _block_refinement_state(tmp_path, utilization=0.45)
    state["block_occupancy_report"] = BlockOccupancyAnalyzer().analyze(state)
    state["block_vlm_review"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "status": "underfilled",
                "severity": "medium",
                "description": "large whitespace remains",
            }
        ]
    }
    before_section = deepcopy(state["story_board"]["spatial_content_plan"]["sections"][0])

    def fake_expansion(self, state, actions, section_by_id):
        return {
            "method": {
                "new_bullets": [
                    "The method uses paper-grounded evidence to add detail while preserving the original section assignment."
                ]
            }
        }

    monkeypatch.setattr(BlockContentRefiner, "_generate_expansion_patches", fake_expansion)
    patch = BlockContentRefiner().refine(state)
    after_section = state["story_board"]["spatial_content_plan"]["sections"][0]

    assert patch["applied"] is True
    assert len(after_section["text_content"]) == len(before_section["text_content"]) + 1
    assert after_section["section_id"] == before_section["section_id"]
    assert after_section["slot_id"] == before_section["slot_id"]
    assert after_section["visual_assets"] == before_section["visual_assets"]


def test_block_content_refiner_reduces_crowded_block_without_changing_refs(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=0.99)
    section = state["story_board"]["spatial_content_plan"]["sections"][0]
    section["text_content"] = [
        "A long factual bullet that is useful but can be shortened when the block is crowded by too much text.",
        "A second factual bullet describing a supporting detail from the paper.",
        "A third lower-priority factual bullet that can be removed first.",
        "A fourth low-priority factual bullet that should be removed in a crowded block.",
    ]
    state["block_occupancy_report"] = BlockOccupancyAnalyzer().analyze(state)
    state["block_vlm_review"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "status": "crowded",
                "severity": "medium",
                "description": "block looks too dense",
            }
        ]
    }
    before_section = deepcopy(section)

    patch = BlockContentRefiner().refine(state)
    after_section = state["story_board"]["spatial_content_plan"]["sections"][0]

    assert patch["applied"] is True
    assert len(after_section["text_content"]) < len(before_section["text_content"])
    assert after_section["section_id"] == before_section["section_id"]
    assert after_section["slot_id"] == before_section["slot_id"]
    assert after_section["visual_assets"] == before_section["visual_assets"]


def test_truncation_removes_dangling_connector_suffixes():
    planner = TemplateBlockPlanner()
    refiner = BlockContentRefiner()

    assert planner._truncate_on_word_boundary(
        "Perturbations include Random, Equal, Flip, and Mixed strategies; a stable method should preserve rankings and identify anomalous annotators.",
        108,
    ).endswith("rankings.")
    assert refiner._truncate_on_word_boundary(
        "A representative arena is Chatbot Arena, where two models answer the same prompt and annotators compare outputs.",
        86,
    ).endswith("prompt.")
    assert refiner._truncate_on_word_boundary(
        "The method identifies anomalous annotators and improves fit and generalization through annotator-aware modeling.",
        81,
    ).endswith("generalization.")
    assert normalize_text_for_poster("Each visit consumes scarce outreach budget and may also.").endswith("budget.")
    assert normalize_text_for_poster("HAGS scales city-wide search by first selecting a.").endswith("search.")
    assert normalize_text_for_poster("The policy then selects the next property within the chosen region using local.").endswith("region.")
    assert normalize_text_for_poster("Updates labels instead of relying only on stale.").endswith("labels.")
    assert normalize_text_for_poster("The predictor is updated with binary cross-entropy, letting.").endswith("entropy.")
    assert normalize_text_for_poster("Teams work with limited staff, limited.").endswith("staff.")
    assert normalize_text_for_poster("judges differ in quality, consistency,.").endswith("consistency.")
    assert normalize_text_for_poster("the method improves generalization through.").endswith("generalization.")
    assert normalize_text_for_poster("Parcels are queried within a budget, while costs may.").endswith("budget.")
    assert normalize_text_for_poster("The search policy is trained with R.").endswith("trained.")
    assert normalize_text_for_poster("Data uses tabular features and imagery; targets are.").endswith("imagery.")
    assert normalize_text_for_poster("Formalizes eviction outreach as a sequential.").endswith("outreach.")
    assert normalize_text_for_poster("HAGS scales this process by splitting.").endswith("process.")
    assert normalize_text_for_poster("HAGS achieves the best target discovery across all budgets, typically.").endswith("budgets.")
    assert normalize_text_for_poster("choose a parcel inside that region..").endswith("region.")


def test_block_content_refiner_forces_min_budget_for_vlm_underfilled_keep(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=0.87)
    state["block_occupancy_report"] = BlockOccupancyAnalyzer().analyze(state)
    block = state["block_occupancy_report"]["blocks"][0]
    block["action"] = "keep"
    block["target_extra_chars"] = 0
    state["block_vlm_review"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "status": "underfilled",
                "severity": "medium",
                "description": "visible whitespace remains",
            }
        ]
    }

    actions = BlockContentRefiner()._decide_actions(state)

    assert actions
    assert actions[0]["action"] == "expand"
    assert actions[0]["target_extra_chars"] > 0


def test_block_content_refiner_compresses_caption_for_visual_too_small(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=0.99)
    state["styled_layout"].append({
        "type": "visual",
        "id": "method_figure",
        "section_id": "method",
        "lane_id": "slot_1",
        "slot_id": "slot_1",
        "x": 1.3,
        "y": 3.0,
        "width": 6.0,
        "height": 3.0,
        "visual_id": "figure_1",
    })
    state["block_occupancy_report"] = BlockOccupancyAnalyzer().analyze(state)
    state["block_vlm_review"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "status": "visual_too_small",
                "severity": "high",
                "description": "figure labels are unreadable",
            }
        ]
    }

    actions = BlockContentRefiner()._decide_actions(state)

    assert actions
    assert actions[0]["action"] == "reduce"
    assert actions[0]["reason"] == "compress text to prioritize visual scale for unreadable figure/table labels"


def test_block_content_refiner_expands_underfilled_visual_too_small_block(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=0.74)
    state["styled_layout"].append({
        "type": "visual",
        "id": "result_table",
        "section_id": "method",
        "lane_id": "slot_1",
        "slot_id": "slot_1",
        "x": 1.3,
        "y": 3.0,
        "width": 6.0,
        "height": 2.0,
        "visual_id": "table_1",
    })
    state["block_occupancy_report"] = BlockOccupancyAnalyzer().analyze(state)
    state["block_vlm_review"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "status": "visual_too_small",
                "severity": "high",
                "description": "table labels are small, but the block has empty lower space",
            }
        ]
    }
    state["block_refinement_count"] = 1

    actions = BlockContentRefiner()._decide_actions(state)

    assert actions
    assert actions[0]["action"] == "expand"
    assert actions[0]["target_extra_chars"] >= 120


def test_block_content_refiner_keeps_medium_crowded_block_below_hard_max(tmp_path):
    state = _block_refinement_state(tmp_path, utilization=0.973)
    state["block_occupancy_report"] = BlockOccupancyAnalyzer().analyze(state)
    block = state["block_occupancy_report"]["blocks"][0]
    block["action"] = "keep"
    block["utilization"] = 0.973
    state["block_vlm_review"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "status": "crowded",
                "severity": "medium",
                "description": "dense but not overflowing",
            }
        ]
    }

    actions = BlockContentRefiner()._decide_actions(state)

    assert actions == []


def test_block_content_refiner_fast_reduce_preserves_min_chars(tmp_path):
    refiner = BlockContentRefiner()
    bullets = [
        "PhishAgent combines online and offline knowledge retrieval with multimodal webpage evidence to address stale brand knowledge, local-brand ambiguity, and delayed phishing indicators.",
        "The agent keeps latency low by using a single-iteration workflow that retrieves candidate brand evidence before making the final phishing decision.",
        "This design targets realistic phishing pages where screenshots, logos, HTML, and domain cues must be interpreted together.",
    ]
    section = {
        "min_chars": 450,
        "max_chars": 550,
        "capacity_budget": {"min_chars": 450, "max_chars": 550},
    }
    action = {"vlm_status": "overflow", "utilization": 0.99}

    reduced = refiner._reduce_bullets_fast(bullets, action, section)

    assert sum(len(item) for item in reduced) >= 450
    assert len(reduced) == len(bullets)


def test_block_content_refiner_fast_skips_medium_crowded_reduce(tmp_path):
    state = create_state(str(tmp_path / "paper.pdf"), enable_block_vlm_review=True)
    state["template_fast_mode"] = True
    state["block_occupancy_report"] = {
        "settings": {"acceptable_min": 0.90, "hard_max": 0.98},
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "motivation",
                "utilization": 0.989,
                "action": "reduce",
                "target_extra_chars": 0,
            }
        ],
    }
    state["block_vlm_review"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "motivation",
                "status": "crowded",
                "severity": "medium",
            }
        ]
    }

    actions = BlockContentRefiner()._decide_actions(state)

    assert actions == []


def test_block_content_refiner_fast_allows_light_text_fill_repair(tmp_path):
    state = create_state(str(tmp_path / "paper.pdf"), enable_block_vlm_review=True)
    state["template_fast_mode"] = True
    state["block_occupancy_report"] = {
        "settings": {"acceptable_min": 0.90, "hard_max": 0.98},
        "blocks": [
            {
                "slot_id": "slot_6",
                "section_id": "results",
                "utilization": 0.89,
                "action": "expand",
                "target_extra_chars": 520,
                "visual_count": 0,
                "reason": "real content leaves bottom whitespace",
            }
        ],
    }
    state["block_vlm_review"] = {"blocks": [{"slot_id": "slot_6", "section_id": "results", "status": "underfilled", "severity": "medium"}]}

    actions = BlockContentRefiner()._decide_actions(state)

    assert actions
    assert actions[0]["action"] == "expand"
    assert actions[0]["target_extra_chars"] <= 420


def test_block_vlm_reviewer_falls_back_when_request_fails(tmp_path, monkeypatch):
    preview_path = tmp_path / "preview.png"
    Image.new("RGB", (900, 500), "white").save(preview_path)
    state = create_state(str(tmp_path / "paper.pdf"), enable_block_vlm_review=True)
    state["output_dir"] = str(tmp_path / "output")
    state["poster_width"] = 9.0
    state["poster_height"] = 5.0
    state["poster_preview_path"] = str(preview_path)
    state["block_occupancy_report"] = {
        "blocks": [
            {
                "slot_id": "slot_1",
                "section_id": "method",
                "section_title": "Method",
                "bbox": {"x": 1.0, "y": 1.0, "w": 4.0, "h": 2.0},
                "utilization": 0.99,
                "action": "reduce",
                "visual_count": 1,
                "reason": "above hard max",
            }
        ]
    }
    monkeypatch.setenv("VLM_BASE_URL", "https://example.invalid/v1")
    monkeypatch.setenv("VLM_API_KEY", "test-key")
    monkeypatch.setenv("VLM_MODEL", "gpt-5.4")
    reviewer = BlockVLMReviewer()
    monkeypatch.setattr(reviewer, "_post_vlm_request", lambda *args, **kwargs: (_ for _ in ()).throw(RuntimeError("bad gateway")))

    result = reviewer(state)

    assert result["errors"] == []
    review = result["block_vlm_review"]
    assert review["source"] == "fallback"
    assert review["blocks"][0]["status"] == "crowded"
    assert review["contact_sheet_path"].endswith("block_contact_sheet.png")
    assert "bad gateway" in review["warnings"][0]


def test_vlm_layout_reviewer_rejects_unresolved_template_whitespace(tmp_path):
    state = create_state(str(tmp_path / "paper.pdf"), enable_vlm_layout_review=True)
    state["output_dir"] = str(tmp_path / "output")
    state["template_layout_mode"] = "template_prior"
    content_dir = Path(state["output_dir"]) / "content"
    content_dir.mkdir(parents=True)
    (content_dir / "micro_layout_report.json").write_text(
        json.dumps(
            {
                "lanes": [
                    {"lane_id": "slot_4", "final_utilization": 0.72},
                    {"lane_id": "slot_1", "final_utilization": 0.94},
                ]
            }
        )
    )
    state["styled_layout"] = [
        {
            "type": "section_container",
            "section_id": "main_results",
            "slot_id": "slot_4",
            "lane_id": "slot_4",
        }
    ]
    review = {
        "overall_score": 90,
        "accept": True,
        "issues": [
            {
                "severity": "medium",
                "category": "whitespace",
                "target": "main_results",
                "description": "large empty area remains",
            }
        ],
        "patch": [],
    }

    gated = VLMLayoutReviewer()._enforce_template_acceptance_gate(review, state)

    assert gated["accept"] is False
    assert any("unresolved_whitespace=True" in warning for warning in gated["warnings"])


def test_vlm_layout_reviewer_accepts_clean_draft_after_max_template_repair(tmp_path):
    state = create_state(str(tmp_path / "paper.pdf"), enable_vlm_layout_review=True)
    state["output_dir"] = str(tmp_path / "output")
    state["template_layout_mode"] = "template_prior"
    state["template_repair_count"] = 1
    content_dir = Path(state["output_dir"]) / "content"
    content_dir.mkdir(parents=True)
    (content_dir / "micro_layout_report.json").write_text(
        json.dumps({"validation": {"issues": []}, "lanes": []})
    )
    review = {
        "overall_score": 72,
        "accept": False,
        "issues": [{"severity": "high", "category": "overflow", "target": "sec_10"}],
        "warnings": [],
    }

    accepted = VLMLayoutReviewer()._accept_after_max_template_repair(review, state)

    assert accepted["accept"] is True
    assert any("accepted after max repair count" in warning for warning in accepted["warnings"])
